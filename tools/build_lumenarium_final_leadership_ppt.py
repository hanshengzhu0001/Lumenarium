from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from tools.build_lumenarium_intern_summary import (
    BLUE, DARK, GREEN, GRID, LIGHT, MID, NAVY, NAVY2, ORANGE, PALE_BLUE,
    PALE_ORANGE, PALE_RED, PALE_TEAL, PURPLE, RED, TEAL, TEAL_DARK, WHITE,
    add_arrow, add_box, add_bullets, add_footer, add_metric, add_picture_cover,
    add_runs, add_stage, add_table, add_text, add_title, set_fill,
)


ROOT = Path(__file__).resolve().parents[1]
VIS = ROOT / "visual_results"
ASSETS = ROOT / "docs" / "assets"
OUT = ROOT / "Lumenarium_暑期实习总结_领导汇报_最终版_2026-08-19.pptx"
NOTES = ROOT / "Lumenarium_暑期实习总结_最终版_逐页讲稿.md"


def base_slide(prs: Presentation, title: str, subtitle: str, page: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, title, subtitle, page, total)
    add_footer(slide, "Lumenarium · 暑期实习总结 · Hansen Zhu · 2026-08")
    return slide


def add_four_up(slide, files, labels, y=1.55, h=3.92):
    xs = [0.55, 3.74, 6.93, 10.12]
    for x, file, label in zip(xs, files, labels):
        add_picture_cover(slide, file, x, y, 2.72, h)
        add_box(slide, x + 0.10, y + h - 0.54, 2.52, 0.42, NAVY, NAVY)
        add_text(slide, label, x + 0.17, y + h - 0.44, 2.38, 0.22,
                 11.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Lumenarium 暑期实习总结"
    prs.core_properties.subject = "从单图重建到更快、可证明、可部署的 3D 场景系统"
    prs.core_properties.author = "Hansen Zhu"
    total = 16

    living = [VIS / "livingroom_10_worst_input.png", VIS / "livingroom_10_worst_v1_final.png",
              VIS / "livingroom_10_worst_v3_final.png", VIS / "livingroom_v5_final.png"]
    office = [VIS / "official_02_worst_input.png", VIS / "official_02_worst_v1_final.png",
              VIS / "official_02_worst_v3_final.png", VIS / "official_v5_final.png"]

    # 1 — cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.13)); set_fill(band, TEAL); band.line.fill.background()
    add_picture_cover(s, living[3], 8.25, 0, 5.08, 7.5, border=NAVY, border_width=0)
    add_text(s, "SUMMER INTERNSHIP REVIEW · 2026", .78, .67, 6.5, .3, 13, TEAL, True)
    add_text(s, "Lumenarium", .78, 1.43, 7.0, .72, 42, WHITE, True)
    add_text(s, "从一张图，到更快、可验证、可交付的 3D 场景", .80, 2.28, 7.20, 1.10, 27, WHITE, True)
    add_box(s, .80, 4.00, 6.36, 1.18, NAVY2, TEAL)
    add_runs(s, [("3.51× S4 提速  ", TEAL, True), ("+7.52 pp physical macro", WHITE, False)], 1.08, 4.31, 5.82, .5, 17, MSO_ANCHOR.MIDDLE)
    add_text(s, "Hansen Zhu · Mentor: Calvin Gu · 腾讯暑期实习", .81, 6.40, 6.3, .28, 13, WHITE, True)

    # 2 — decision frame
    s = base_slide(prs, "先定义“好场景”：光子真正需要什么？", "不是所有指标同等重要；两个月的所有技术选择都围绕这个排序", 2, total)
    add_box(s, .62, 1.40, 7.22, 4.95, WHITE, TEAL, width=1.7)
    add_text(s, "交付优先级", .95, 1.72, 2.5, .36, 21, TEAL_DARK, True)
    priorities = [("01  资产完整性", "主体与关键道具不能缺席"), ("02  速度", "能进入真实制作链路，而不是离线实验"),
                  ("03  不穿模", "最影响可信度的三维错误"), ("04  物体间关系", "桌上、墙上、床上等结构要成立")]
    for i, (head, body) in enumerate(priorities):
        y = 2.31 + i * .82
        add_text(s, head, 1.00, y, 2.62, .30, 16, NAVY, True)
        add_text(s, body, 3.52, y, 3.74, .31, 14, DARK)
    add_box(s, 8.15, 1.40, 4.56, 4.95, PALE_ORANGE, ORANGE, width=1.5)
    add_text(s, "可以接受的 trade-off", 8.48, 1.72, 3.60, .36, 20, ORANGE, True)
    add_bullets(s, [("单体资产 pose：", "rotation / translation 可略有牺牲。"),
                    ("少量非关键悬空：", "优先显式 unresolved，不伪装成功。")],
                8.48, 2.42, 3.48, 1.54, 14, bullet_color=ORANGE)
    add_box(s, 8.48, 4.44, 3.50, 1.20, NAVY, NAVY)
    add_text(s, "核心判断", 8.73, 4.66, 1.2, .26, 14, TEAL, True)
    add_text(s, "场景整体可用性 > 单个资产像素级复刻", 8.73, 5.07, 2.98, .42, 14, WHITE, True)

    # 3 — challenge
    s = base_slide(prs, "原始系统不是“不能生成”，而是难以稳定交付", "看起来像一个房间，不等于三维关系真的成立", 3, total)
    add_picture_cover(s, living[1], .62, 1.38, 6.05, 4.90)
    add_box(s, 6.94, 1.38, 5.76, 4.90, WHITE, GRID)
    add_text(s, "四类阻碍", 7.27, 1.72, 2.0, .34, 20, NAVY, True)
    add_bullets(s, [("不完整：", "窗帘、支撑物或关键对象在链路中丢失"), ("太慢：", "S4 的 SA-5000 单场景约 678 秒"),
                    ("会穿模：", "局部视觉合理，但真实网格已重叠"), ("关系不可信：", "错误 parent、悬空和错误上墙会级联")],
                7.28, 2.25, 4.70, 2.60, 14.0, bullet_color=RED)
    add_box(s, 7.24, 5.20, 4.86, .78, PALE_TEAL, TEAL)
    add_text(s, "因此目标不是“多搜索”，而是“只提交能证明的修改”", 7.46, 5.37, 4.44, .46, 12.4, TEAL_DARK, True, PP_ALIGN.CENTER)

    # 4 — journey
    s = base_slide(prs, "两个月的主线：结构 → 速度 → 证明 → 服务", "每一个失败样例都变成了下一层系统能力", 4, total)
    milestones = [("V1", "Baseline", "Imaginarium 重建", BLUE), ("V3", "Support tree", "支撑 / stack", ORANGE),
                  ("V4", "DeepSearch", "更快检索", PURPLE), ("V5", "LM + Proof", "更快、更物理", TEAL),
                  ("API", "Service", "双 A10 交付", GREEN)]
    xs = [.58, 3.12, 5.66, 8.20, 10.74]
    for i, (code, title, detail, color) in enumerate(milestones):
        add_stage(s, xs[i], 2.08, 2.03, 2.35, code, title, detail, color)
        if i < 4: add_arrow(s, xs[i] + 2.06, 3.22, xs[i + 1] - .06, 3.22, MID, 1.8)
    add_box(s, .78, 5.08, 11.76, 1.10, NAVY, NAVY)
    add_text(s, "一句话：从“全局盲搜一个看起来不错的答案”，转向“关系定向修改 + 几何证书提交”。", 1.13, 5.39, 11.05, .38, 18, WHITE, True, PP_ALIGN.CENTER)

    # 5 — workflow
    s = base_slide(prs, "当前 Lumenarium：一张图走完整 S0–S4", "继承 Imaginarium 的视觉重建基础，重做检索、优化、证明和部署", 5, total)
    stages = [("S0", "Geometry", "深度 / 相机", BLUE), ("S1", "Understand", "SAM3 / 场景图", TEAL),
              ("S2", "Retrieve", "DeepSearch / DINOv2", ORANGE), ("S3", "Recover", "pose / support tree", PURPLE),
              ("S4", "Optimize", "SceneLM / SceneProof", GREEN)]
    xs = [.56, 3.09, 5.62, 8.15, 10.68]
    for i, item in enumerate(stages):
        add_stage(s, xs[i], 1.83, 2.08, 2.35, *item)
        if i < 4: add_arrow(s, xs[i] + 2.10, 3.00, xs[i + 1] - .05, 3.00, MID, 1.8)
    add_box(s, .70, 4.72, 12.00, 1.34, NAVY, NAVY)
    add_runs(s, [("SceneLM  ", TEAL, True), ("关系范围内 400-step 优化", WHITE, False), ("     SceneProof  ", TEAL, True), ("证书 → commit / restore / unresolved", WHITE, False)], .98, 5.00, 11.42, .44, 15.5, MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    add_text(s, "输出：editable placement · render · evaluation certificate · result bundle", 1.05, 5.57, 11.30, .25, 11, RGBColor(199, 219, 238), False, PP_ALIGN.CENTER)

    # 6 — living comparison
    s = base_slide(prs, "同一客厅：从“重建出来”到“整体可用”", "Input / Imaginarium V1 / support-aware V3 / final Lumenarium V5", 6, total)
    add_four_up(s, living, ["Input", "V1 · Imaginarium", "V3 · Support-aware", "V5 · Lumenarium"])
    add_box(s, .70, 5.82, 11.96, .58, PALE_TEAL, TEAL)
    add_text(s, "V5 更完整地恢复大件与软装，并把系统目标从单体 pose 转向场景级可用性。", .96, 5.98, 11.42, .28, 14.2, TEAL_DARK, True, PP_ALIGN.CENTER)

    # 7 — office comparison
    s = base_slide(prs, "同一办公室：更快检索，同时保持关键资产覆盖", "视觉结果不是 GT；它展示的是产品级完整性、关系和冲突 trade-off", 7, total)
    add_four_up(s, office, ["Input", "V1 · Imaginarium", "V3 · Support-aware", "V5 · Lumenarium"])
    add_box(s, .70, 5.82, 11.96, .58, PALE_ORANGE, ORANGE)
    add_text(s, "DeepSearch 提升速度和资产覆盖，但单体 rotation / translation 仍需后续校准。", .96, 5.98, 11.42, .28, 14.2, DARK, True, PP_ALIGN.CENTER)

    # 8 — quality
    s = base_slide(prs, "质量结果：V3 找回更多，V5 让场景更物理", "Paper30 · Primary objects · visible mask ≥ 8,000 px · GT 只用于评测", 8, total)
    data = [["版本", "Primary recovery", "Primary parent", "Physical macro", "它解决什么"],
            ["Imaginarium V1", "89.49%", "89.32%", "52.98%", "原始基线"],
            ["Lumenarium V3", "91.40%", "87.80%", "52.14%", "支撑感知恢复"],
            ["Lumenarium V5", "88.22%", "80.14%", "62.10%", "速度 + 物理主版本"]]
    add_table(s, data, .62, 1.43, 12.08, 2.50, [2.4, 1.8, 1.7, 1.65, 2.4], 12.4, {2: PALE_ORANGE, 3: PALE_TEAL})
    add_metric(s, .68, 4.35, 3.70, 1.55, "91.40%", "V3 recovery", "结构化支撑帮助找回 Primary 对象", ORANGE)
    add_metric(s, 4.78, 4.35, 3.70, 1.55, "+7.52 pp", "V5 physical macro", "54.58% V4 → 62.10% V5", GREEN)
    add_box(s, 8.88, 4.35, 3.78, 1.55, PALE_BLUE, BLUE)
    add_text(s, "诚实边界", 9.17, 4.64, 2.0, .30, 17, BLUE, True)
    add_text(s, "DeepSearch 的 pose trade-off 不归因给 SceneLM，也不从报告中隐藏。", 9.17, 5.08, 3.12, .62, 12.2, DARK)

    # 9 — speed
    s = base_slide(prs, "速度结果：V5 把全链路从约 23.8 分钟降到 13.83 分钟", "V1 为历史恢复估算；V4 为 S0–S3 实测 + legacy S4；V5 为 Paper30 冷启动实测", 9, total)
    data = [["版本", "端到端 / scene", "S4 / scene", "相对 V1", "主要变化"],
            ["Imaginarium V1", "≈23.8 min", "677.8 s", "1.00×", "legacy retrieval + SA-5000"],
            ["V4 DeepSearch", "≈21.9 min", "677.8 s", "≈1.09×", "检索提速"],
            ["Lumenarium V5-fast", "13.83 min", "192.9 s", "≈1.72×", "SceneLM + SceneProof"]]
    add_table(s, data, .62, 1.43, 12.08, 2.38, [2.30, 1.72, 1.55, 1.30, 3.03], 12.4, {3: PALE_TEAL})
    add_metric(s, .70, 4.24, 3.72, 1.62, "3.513×", "S4 speedup", "677.770s → 192.930s", TEAL)
    add_metric(s, 4.80, 4.24, 3.72, 1.62, "13.83 min", "V5 S0–S4", "829.879s / scene measured", BLUE)
    add_box(s, 8.90, 4.24, 3.72, 1.62, PALE_ORANGE, ORANGE)
    add_text(s, "下一瓶颈", 9.19, 4.54, 2.0, .30, 17, ORANGE, True)
    add_text(s, "S1 = 443.0s；Gemini 并发与请求合并是下一轮主要空间。", 9.19, 4.99, 3.02, .63, 12.3, DARK)

    # 10 — SceneLM
    s = base_slide(prs, "为什么 SceneLM 更快：从“全屋盲搜”到“关系定向求解”", "LM 不直接决定最终结果；它只缩小搜索空间并提出结构化候选", 10, total)
    add_box(s, .62, 1.43, 4.04, 4.85, PALE_RED, RED)
    add_text(s, "Imaginarium S4", .95, 1.76, 2.4, .34, 20, RED, True)
    add_text(s, "5,000-step simulated annealing", .95, 2.40, 3.15, .36, 18, NAVY, True)
    add_bullets(s, ["全局扰动大量无关对象", "每步反复计算全局能量", "搜索方向弱、停止条件慢"], .96, 3.10, 3.10, 1.70, 14, bullet_color=RED)
    add_box(s, 4.95, 1.43, 7.75, 4.85, PALE_TEAL, TEAL)
    add_text(s, "Lumenarium SceneLM", 5.28, 1.76, 3.3, .34, 20, TEAL_DARK, True)
    add_text(s, "① Relation Program 编译", 5.30, 2.36, 3.0, .31, 15, NAVY, True)
    add_text(s, "support · collision · plane · semantic", 5.31, 2.76, 4.6, .29, 13, MID, False, font="Consolas")
    add_text(s, "② 只更新 implicated objects / DoF", 5.30, 3.30, 4.2, .31, 15, NAVY, True)
    add_text(s, "③ exact leaf-translation Schur elimination", 5.30, 3.87, 5.0, .31, 15, NAVY, True)
    add_box(s, 5.30, 4.58, 6.78, .96, NAVY, NAVY)
    add_text(s, "global SA-5000  →  scoped 400-step optimization", 5.59, 4.87, 6.20, .30, 17, WHITE, True, PP_ALIGN.CENTER, font="Consolas")
    add_text(s, "减少的是无效搜索，不是省掉必要的几何检查。", 5.50, 5.78, 6.35, .27, 13, TEAL_DARK, True, PP_ALIGN.CENTER)

    # 11 — SceneProof
    s = base_slide(prs, "为什么 SceneProof 更物理：proposal 不等于 commit", "每个候选都是一笔事务；只有局部通过且全局非劣才提交", 11, total)
    steps = [("候选", "LM / rule", BLUE), ("重建", "true mesh", PURPLE), ("局部门", "contact / COM", ORANGE),
             ("全局门", "family non-regression", TEAL), ("裁决", "commit / restore", GREEN)]
    xs = [.55, 3.03, 5.51, 7.99, 10.47]
    for i, (title, detail, color) in enumerate(steps):
        add_box(s, xs[i], 1.52, 2.12, 1.18, WHITE, color, width=1.5)
        add_text(s, title, xs[i] + .18, 1.75, 1.76, .30, 16, color, True, PP_ALIGN.CENTER)
        add_text(s, detail, xs[i] + .15, 2.18, 1.82, .25, 10.4, MID, False, PP_ALIGN.CENTER)
        if i < 4: add_arrow(s, xs[i] + 2.15, 2.12, xs[i + 1] - .05, 2.12, MID, 1.8)
    add_box(s, .68, 3.20, 5.80, 2.54, WHITE, GRID)
    add_text(s, "必须出现的证据", .98, 3.51, 2.6, .34, 18, NAVY, True)
    add_bullets(s, ["真实网格不增加碰撞 / 穿透", "declared support contact + COM / boundary", "plane / semantic / visibility 不明显退化"], .98, 4.05, 4.92, 1.30, 13.2)
    add_box(s, 6.78, 3.20, 5.88, 2.54, PALE_TEAL, TEAL)
    add_text(s, "提交规则", 7.10, 3.51, 2.0, .34, 18, TEAL_DARK, True)
    add_text(s, "Accept(Δ) ⇔ LocalGates(Δ) ∧ ΔFamily ≥ −ε", 7.11, 4.11, 5.15, .48, 19, NAVY, True, PP_ALIGN.CENTER, font="Cambria Math")
    add_text(s, "失败 → restore incumbent\n证据不足 → unresolved（不伪装成功）", 7.12, 4.79, 5.12, .68, 14, DARK, True, PP_ALIGN.CENTER)

    # 12 — causal proof
    s = base_slide(prs, "因果闭环：速度与物理提升来自不同机制", "冻结上游后做分层对比，避免把所有收益都归给同一个模块", 12, total)
    add_box(s, .62, 1.42, 5.80, 4.92, PALE_BLUE, BLUE)
    add_text(s, "速度归因", .95, 1.75, 2.0, .34, 20, BLUE, True)
    add_text(s, "V1 → V4", .97, 2.40, 1.4, .31, 16, NAVY, True); add_text(s, "DeepSearch：≈23.8 → ≈21.9 min", 2.24, 2.40, 3.45, .31, 14, DARK)
    add_text(s, "V4 → V5", .97, 3.09, 1.4, .31, 16, NAVY, True); add_text(s, "SceneLM：S4 677.8 → 192.9s", 2.24, 3.09, 3.45, .31, 14, DARK)
    add_box(s, .96, 4.10, 4.96, 1.34, NAVY, NAVY)
    add_text(s, "结论", 1.22, 4.34, 1.0, .27, 14, TEAL, True)
    add_text(s, "DeepSearch 加速 retrieval；SceneLM 加速 optimization。", 1.22, 4.78, 4.34, .42, 14, WHITE, True)
    add_box(s, 6.78, 1.42, 5.88, 4.92, PALE_TEAL, TEAL)
    add_text(s, "物理归因", 7.11, 1.75, 2.0, .34, 20, TEAL_DARK, True)
    add_text(s, "V4 recovery / parent", 7.13, 2.42, 2.42, .30, 15, NAVY, True); add_text(s, "88.22% / 80.14%", 10.05, 2.42, 2.05, .30, 14, DARK)
    add_text(s, "V5 recovery / parent", 7.13, 3.10, 2.42, .30, 15, NAVY, True); add_text(s, "88.22% / 80.14%", 10.05, 3.10, 2.05, .30, 14, DARK)
    add_box(s, 7.12, 4.10, 5.02, 1.34, NAVY, NAVY)
    add_text(s, "同一上游工作点", 7.39, 4.34, 2.2, .27, 14, TEAL, True)
    add_text(s, "physical 54.58% → 62.10%：收益来自 SceneLM + Proof。", 7.39, 4.78, 4.40, .42, 13.5, WHITE, True)

    # 13 — service
    s = base_slide(prs, "不止是论文代码：已经变成双 A10 端到端服务", "新图完整运行 S0–S4；相同输入跨 Fast / Medium / Best 复用冻结缓存", 13, total)
    add_box(s, .62, 1.43, 3.14, 4.82, WHITE, GRID)
    add_text(s, "技术美术入口", .95, 1.76, 2.3, .34, 19, NAVY, True)
    add_bullets(s, ["任意尺寸 PNG / JPEG", "S0–S4 实时进度", "可编辑 placement + render", "证书 + unresolved + ZIP"], .96, 2.38, 2.35, 2.40, 14)
    add_box(s, 4.05, 1.43, 4.18, 4.82, PALE_BLUE, BLUE)
    add_text(s, "可靠调度", 4.38, 1.76, 2.0, .34, 19, BLUE, True)
    add_bullets(s, [("双 A10：", "原子 claim，避免重复领取"), ("恢复：", "回收死亡 claim；单场失败不中断"),
                    ("重试：", "HTTP / S4 缺失有限重试"), ("缓存：", "复用 S0–S3 / Fix61")], 4.39, 2.35, 3.30, 2.78, 12.2, bullet_color=BLUE)
    add_box(s, 8.53, 1.43, 4.14, 4.82, PALE_TEAL, TEAL)
    add_text(s, "三种模式", 8.86, 1.76, 2.0, .34, 19, TEAL_DARK, True)
    add_bullets(s, [("Fast：", "冻结 Fix61，论文定量"), ("Medium：", "visual-safe 展示清理"),
                    ("Best：", "真实支撑审计与 first-contact repair")], 8.87, 2.35, 3.25, 2.48, 12.2, bullet_color=TEAL)
    add_box(s, 8.86, 5.05, 3.25, .56, NAVY, NAVY)
    add_text(s, "embedding.lightart.qq.com", 9.00, 5.21, 2.97, .24, 11.5, WHITE, True, PP_ALIGN.CENTER)

    # 14 — limits
    s = base_slide(prs, "我们知道哪里还不够好，也知道为什么", "可解释的 trade-off 比一个“全都更好”的故事更可信", 14, total)
    add_box(s, .62, 1.43, 5.82, 4.88, PALE_ORANGE, ORANGE)
    add_text(s, "当前限制", .95, 1.77, 2.2, .34, 20, ORANGE, True)
    add_bullets(s, [("Pose：", "DeepSearch 后单体 pose 工作点下降"), ("Support：", "证据不足时 unresolved，不强行修"),
                    ("Latency：", "S1 Gemini 场景图与语义 API 仍占 443s")], .96, 2.37, 4.92, 2.42, 12.6, bullet_color=ORANGE)
    add_text(s, "这不是被隐藏的失败，而是下一步研究边界。", .98, 5.45, 4.86, .32, 14, DARK, True)
    add_box(s, 6.78, 1.43, 5.88, 4.88, PALE_TEAL, TEAL)
    add_text(s, "下一步最高 ROI", 7.11, 1.77, 2.8, .34, 20, TEAL_DARK, True)
    add_bullets(s, [("S1 并发：", "Gemini 8 并发 + 合并请求，目标 250–320s"), ("Pose 校准：", "seed-locked S2 ablation + frame 校正"),
                    ("DCC 闭环：", "UE / Blender 导入、人工确认与回写")], 7.12, 2.37, 4.95, 2.42, 12.6, bullet_color=TEAL)
    add_box(s, 7.11, 5.10, 5.04, .64, NAVY, NAVY)
    add_text(s, "原则：先攻最大瓶颈，再增加模型复杂度。", 7.37, 5.30, 4.53, .27, 13.2, WHITE, True, PP_ALIGN.CENTER)

    # 15 — impact
    s = base_slide(prs, "这两个月留下的，不只是一个更高的分数", "算法、证据、工具链和产品入口同时闭环", 15, total)
    cards = [("研究", "两项方法创新", "support-aware reconstruction\nSceneLM + SceneProof", BLUE),
             ("工程", "可复现评测", "Paper30 / 8000px+\nprovenance / rollback", ORANGE),
             ("产品", "双 A10 服务", "Web / API / cache\nFast / Medium / Best", TEAL)]
    for i, (tag, head, body, color) in enumerate(cards):
        x = .70 + i * 4.18
        add_box(s, x, 1.55, 3.78, 3.50, WHITE, color, width=1.7)
        add_text(s, tag, x + .30, 1.88, 1.0, .28, 13, color, True)
        add_text(s, head, x + .30, 2.41, 3.02, .38, 21, NAVY, True)
        add_text(s, body, x + .30, 3.17, 3.02, 1.05, 15, DARK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_box(s, 1.10, 5.48, 11.14, .72, NAVY, NAVY)
    add_text(s, "把“单图 3D 场景生成”从研究原型推进成：更快、可证明、能被技术美术直接使用的系统。", 1.38, 5.70, 10.58, .31, 17, WHITE, True, PP_ALIGN.CENTER)

    # 16 — close
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.13)); set_fill(band, TEAL); band.line.fill.background()
    add_text(s, "Lumenarium", .84, .72, 4.6, .45, 27, TEAL, True)
    add_text(s, "我没有让系统“搜索得更多”", .84, 1.63, 10.5, .55, 30, WHITE, True)
    add_text(s, "而是让它只修改该修改的，并证明每一次提交。", .84, 2.43, 11.3, .60, 29, WHITE, True)
    add_box(s, .86, 3.58, 11.60, 1.36, NAVY2, TEAL)
    add_runs(s, [("完整性  ", TEAL, True), ("→ 速度  ", WHITE, False), ("→ 低穿模  ", WHITE, False),
                 ("→ 关系可信  ", WHITE, False), ("→ 可交付", WHITE, False)], 1.20, 3.98, 10.90, .48, 20, MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    add_text(s, "谢谢 · Questions", .86, 5.83, 4.0, .42, 24, WHITE, True)
    add_text(s, "Hansen Zhu · Mentor: Calvin Gu", .88, 6.42, 5.0, .28, 13, RGBColor(202, 220, 239))
    add_text(s, "Demo: bilibili.com/video/BV1tpbD6hERB", 7.00, 6.43, 5.45, .25, 10.5, RGBColor(202, 220, 239), False, PP_ALIGN.RIGHT)

    prs.save(OUT)
    return OUT


def write_notes() -> Path:
    text = "# Lumenarium 暑期实习总结：逐页讲稿\n\n"
    text += "1. **封面**：两个月的核心不是堆功能，而是把原型推进成更快、可证明、可交付的系统。\n"
    text += "2. **优先级**：先把光子真实关心的四件事讲清楚，后面的所有取舍都有依据。\n"
    text += "3. **问题**：原系统能生成，但完整性、速度、穿模与关系错误阻碍交付。\n"
    text += "4. **路径**：V1 到 V3 建结构，V4 换检索，V5 重做优化与证明，最后产品化。\n"
    text += "5. **全链路**：S0–S3 负责理解与恢复，S4 负责优化与证书。\n"
    text += "6–7. **视觉对比**：强调相同输入和系统级可用性，不宣称单体 pose 全面更优。\n"
    text += "8. **质量**：V3 recovery 最好；V5 physical macro 最好，说明两阶段解决不同问题。\n"
    text += "9. **速度**：V4 省 retrieval，V5 省 optimization；全链路约 23.8 → 13.83 分钟。\n"
    text += "10. **SceneLM**：把全局 SA-5000 换成关系范围内 400-step 优化与 Schur elimination。\n"
    text += "11. **SceneProof**：模型只提 proposal，几何证书决定 commit、restore 或 unresolved。\n"
    text += "12. **因果归因**：冻结上游后，速度和 physical 的提升都能对应到具体模块。\n"
    text += "13. **服务**：双 A10、Web/API、缓存和恢复机制让技术美术能直接使用。\n"
    text += "14. **限制**：主动讲 pose、unresolved 和 S1 latency，体现研究判断而不是包装。\n"
    text += "15. **影响**：方法、评测、工程、产品四条线均形成可复用资产。\n"
    text += "16. **收尾**：不是让系统搜索更多，而是只修改该修改的，并证明每一次提交。\n"
    NOTES.write_text(text, encoding="utf-8")
    return NOTES


if __name__ == "__main__":
    print(build())
    print(write_notes())
