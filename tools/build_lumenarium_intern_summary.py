from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT.parent / "Imaginarium_v1_v3_Professional.pptx"
OUT = ROOT / "Lumenarium_暑期实习总结_领导汇报_2026-08-17.pptx"
ASSETS = ROOT / "intern_summary_assets"

NAVY = RGBColor(18, 38, 95)
NAVY2 = RGBColor(28, 62, 116)
TEAL = RGBColor(0, 196, 157)
TEAL_DARK = RGBColor(0, 142, 118)
BLUE = RGBColor(48, 105, 210)
ORANGE = RGBColor(242, 151, 40)
RED = RGBColor(221, 73, 73)
GREEN = RGBColor(35, 158, 96)
PURPLE = RGBColor(117, 85, 190)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(245, 247, 251)
LIGHT2 = RGBColor(235, 240, 248)
GRID = RGBColor(213, 220, 232)
MID = RGBColor(99, 111, 132)
DARK = RGBColor(29, 45, 88)
PALE_TEAL = RGBColor(226, 247, 241)
PALE_BLUE = RGBColor(232, 239, 252)
PALE_ORANGE = RGBColor(255, 244, 226)
PALE_RED = RGBColor(255, 237, 237)
FONT = "Microsoft YaHei"
MONO = "Consolas"


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def set_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(slide, text: str, x: float, y: float, w: float, h: float,
             size: float = 20, color: RGBColor = DARK, bold: bool = False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font: str = FONT,
             margin: float = 0.02, italic: bool = False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box


def add_runs(slide, runs, x, y, w, h, size=18, valign=MSO_ANCHOR.TOP,
             align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.03)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    for text, color, bold in runs:
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return box


def add_box(slide, x, y, w, h, fill=WHITE, line=GRID, radius=True, width=1.1):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill); set_line(shape, line, width)
    return shape


def add_title(slide, title: str, subtitle: str, page: int, total: int):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(0.42), Inches(0.16), Inches(0.64))
    set_fill(bar, TEAL); bar.line.fill.background()
    add_text(slide, title, 0.90, 0.37, 10.7, 0.52, 26, NAVY, True)
    add_text(slide, subtitle, 0.90, 0.89, 10.9, 0.27, 11.5, MID)
    add_text(slide, f"{page:02d} / {total:02d}", 11.75, 0.48, 0.98, 0.25, 10.5, MID, True, PP_ALIGN.RIGHT)


def add_footer(slide, text="Lumenarium · 暑期实习总结 · Hansen Zhu · 2026-08"):
    add_text(slide, text, 0.64, 7.22, 12.0, 0.17, 8, MID)


def add_bullets(slide, items: Iterable, x, y, w, h, size=16, color=DARK,
                bullet_color=TEAL, spacing=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        head, body = item if isinstance(item, tuple) else ("", item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing); p.line_spacing = 1.08
        r = p.add_run(); r.text = "▌ "; r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = bullet_color
        if head:
            r = p.add_run(); r.text = head; r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
        r = p.add_run(); r.text = body; r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
    return box


def add_metric(slide, x, y, w, h, value, label, note, accent=TEAL):
    add_box(slide, x, y, w, h, WHITE, GRID)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.11), Inches(h))
    set_fill(band, accent); band.line.fill.background()
    add_text(slide, value, x + 0.25, y + 0.18, w - 0.38, 0.48, 27, accent, True)
    add_text(slide, label, x + 0.25, y + 0.69, w - 0.38, 0.31, 12.5, DARK, True)
    add_text(slide, note, x + 0.25, y + 1.05, w - 0.38, h - 1.15, 9.5, MID)


def add_picture_cover(slide, path: Path, x, y, w, h, border=GRID, border_width=1.0):
    with Image.open(path) as im:
        iw, ih = im.size
    target, source = w / h, iw / ih
    crop_l = crop_r = crop_t = crop_b = 0.0
    if source > target:
        visible = target / source; crop_l = crop_r = (1 - visible) / 2
    else:
        visible = source / target; crop_t = crop_b = (1 - visible) / 2
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    pic.crop_left = crop_l; pic.crop_right = crop_r; pic.crop_top = crop_t; pic.crop_bottom = crop_b
    set_line(pic, border, border_width)
    return pic


def add_table(slide, data, x, y, w, h, col_widths=None, font_size=12,
              highlight_rows=None):
    rows, cols = len(data), len(data[0])
    tshape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = tshape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(Inches(w) * cw / total)
    highlight_rows = highlight_rows or {}
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY2 if r == 0 else highlight_rows.get(r, WHITE if r % 2 else LIGHT)
            cell.text = str(data[r][c])
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            p.vertical_anchor = MSO_ANCHOR.MIDDLE
            for run in p.runs:
                run.font.name = FONT; run.font.size = Pt(font_size); run.font.bold = (r == 0 or c == 0)
                run.font.color.rgb = WHITE if r == 0 else DARK
    return tshape


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, width=2.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    set_line(line, color, width)
    line.line.end_arrowhead = True
    return line


def add_stage(slide, x, y, w, h, code, title, detail, accent):
    add_box(slide, x, y, w, h, WHITE, accent, width=1.5)
    add_text(slide, code, x + 0.18, y + 0.16, 0.48, 0.30, 16, accent, True)
    add_text(slide, title, x + 0.18, y + 0.55, w - 0.36, 0.31, 15, DARK, True)
    add_text(slide, detail, x + 0.18, y + 0.98, w - 0.36, h - 1.10, 10.2, MID)


def build_speed_chart() -> Path:
    ASSETS.mkdir(exist_ok=True)
    out = ASSETS / "speed_breakdown.png"
    W, H = 1500, 620
    im = Image.new("RGB", (W, H), "white"); d = ImageDraw.Draw(im)
    font_path = ROOT / "utils" / "SimHei.ttf"
    f22 = ImageFont.truetype(str(font_path), 34); f16 = ImageFont.truetype(str(font_path), 25); f13 = ImageFont.truetype(str(font_path), 20)
    values = [("S0", 9.687, (48,105,210)), ("S1",443.036,(0,196,157)), ("S2",137.451,(242,151,40)), ("S3",44.790,(117,85,190)), ("调度",1.986,(150,158,175))]
    total = sum(v for _, v, _ in values); x0, x1, y0, bh = 105, 1410, 145, 95
    d.text((105, 42), "冷启动 S0–S3：636.949 秒 / 场景", fill=(18,38,95), font=f22)
    cur = x0
    for name, val, color in values:
        bw = (x1-x0) * val / total
        d.rectangle((cur,y0,cur+bw,y0+bh), fill=color)
        if bw > 70: d.text((cur+10,y0+28), name, fill="white", font=f16)
        cur += bw
    legend_y = 290
    for i, (name,val,color) in enumerate(values):
        xx = 105 + (i % 3) * 430; yy = legend_y + (i // 3) * 72
        d.rounded_rectangle((xx,yy,xx+34,yy+34), 6, fill=color)
        d.text((xx+50,yy-1), f"{name}  {val:.1f}s  ({val/total*100:.1f}%)", fill=(29,45,88), font=f16)
    d.text((105, 475), "结论：S1 占 69.6%，是下一阶段端到端提速的主瓶颈；S4 已完成结构性加速。", fill=(99,111,132), font=f13)
    im.save(out)
    return out


def new_slide(prs, title, subtitle, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, title, subtitle, page, total)
    add_footer(slide)
    return slide


def build() -> Path:
    speed_chart = build_speed_chart()
    office = ROOT / "docs/assets/lumenarium_a10_office_demo.png"
    living = ROOT / "docs/assets/lumenarium_a10_livingroom_demo.png"
    pipeline = ROOT / "media/pipeline.png"
    before = ROOT / "fix116_final.png"
    after = ROOT / "visual_safe_runnerfix2_98323a90.png"

    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)
    prs.core_properties.title = "Lumenarium 暑期实习总结"
    prs.core_properties.subject = "单图到可编辑三维场景：SceneLM、SceneProof 与双 A10 服务"
    prs.core_properties.author = "Hansen Zhu"
    total = 18

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14)); set_fill(accent, TEAL); accent.line.fill.background()
    poly = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(8.15), Inches(-0.15), Inches(5.5), Inches(7.7)); set_fill(poly, NAVY2, 18); poly.line.fill.background()
    add_text(s, "SUMMER INTERNSHIP REVIEW · 2026", 0.86, 0.66, 7.8, 0.28, 14, TEAL, True)
    add_text(s, "Lumenarium", 0.86, 1.40, 7.8, 0.68, 42, WHITE, True)
    add_text(s, "从一张图，到可编辑、可验证的 3D 场景", 0.89, 2.23, 10.6, 0.62, 28, WHITE, True)
    add_text(s, "两个月：把研究原型推进为可复现算法、论文实验与在线服务", 0.90, 3.02, 9.8, 0.36, 17, RGBColor(205,221,247))
    add_box(s, 0.90, 4.08, 6.35, 1.14, NAVY2, TEAL)
    add_runs(s, [("核心结果  ", TEAL, True), ("Physical macro +7.52 pp  ·  S4 3.51×  ·  双 A10 全链路部署", WHITE, False)], 1.12, 4.34, 5.9, 0.62, 15.5, MSO_ANCHOR.MIDDLE)
    add_text(s, "Hansen Zhu  ·  Mentor: Calvin Gu  ·  腾讯暑期实习", 0.91, 6.31, 7.8, 0.30, 14, WHITE, True)
    add_text(s, "Bilibili Demo: BV1tpbD6hERB", 0.91, 6.72, 5.3, 0.26, 11, RGBColor(205,221,247))

    # 2 Executive summary
    s = new_slide(prs, "先讲结论：两个月，完成三次跨越", "从“能生成”到“更合理、更快、可交付”", 2, total)
    add_metric(s, 0.62, 1.42, 3.82, 1.55, "+7.52 pp", "物理质量提升", "V4 DeepSearch 54.58% → V5-fast 62.10%", GREEN)
    add_metric(s, 4.75, 1.42, 3.82, 1.55, "3.51×", "S4 优化提速", "Legacy SA-5000 677.8s → SceneLM/Fix61 192.9s", TEAL)
    add_metric(s, 8.88, 1.42, 3.82, 1.55, "S0–S4", "端到端产品化", "双 A10、Web/API、缓存、证书、结果包", BLUE)
    add_box(s, 0.62, 3.32, 12.08, 2.88, WHITE, GRID)
    add_text(s, "一句话概括", 0.93, 3.63, 2.0, 0.32, 18, NAVY, True)
    add_text(s, "我没有继续堆更多搜索步数，而是让模型只改“出问题的关系”，再让几何证书决定是否提交。", 0.93, 4.10, 11.15, 0.58, 23, DARK, True)
    add_bullets(s, [("SceneLM：", "把全局盲搜变成关系范围内的定向优化。"), ("SceneProof：", "每次修改都必须携带碰撞、支撑和非退化证据。"), ("Lumenarium Service：", "把研究链路封装成技术美术可直接上传图片使用的服务。")], 0.96, 4.91, 11.15, 1.10, 14)

    # 3 Problem story
    s = new_slide(prs, "故事从一个反直觉的问题开始", "“看起来像”不等于“在三维里真的成立”", 3, total)
    add_picture_cover(s, before, 0.62, 1.40, 6.02, 4.85)
    add_box(s, 6.92, 1.40, 5.78, 4.85, WHITE, GRID)
    add_text(s, "原始系统已经能重建场景，但交付时仍会出戏", 7.25, 1.73, 5.08, 0.62, 21, NAVY, True)
    add_bullets(s, [("悬空：", "物体位于“正确高度”，却没有真实支撑。"), ("穿模：", "局部看似合理，三维网格已相交。"), ("错误父子关系：", "窗户跟随画框、物体被吸向错误墙面。"), ("速度：", "S4 需要 5,000 步模拟退火，单场景约 678 秒。")], 7.25, 2.63, 4.96, 2.45, 15)
    add_box(s, 7.23, 5.38, 4.96, 0.55, PALE_ORANGE, ORANGE)
    add_text(s, "目标：不是“多修一点”，而是“只提交可证明的改动”。", 7.43, 5.53, 4.57, 0.24, 13.5, DARK, True)

    # 4 Pipeline
    s = new_slide(prs, "Lumenarium 全链路：五个阶段，一份可编辑结果", "新图片完整运行 S0–S4；相同图片可跨模式复用冻结缓存", 4, total)
    xs = [0.58, 3.10, 5.62, 8.14, 10.66]
    stages = [("S0","Geometry","深度、相机、房间几何",BLUE),("S1","Parsing","检测、分割、场景图",TEAL),("S2","Retrieval","DeepSearch 资产检索",ORANGE),("S3","Pose","姿态与堆叠关系",PURPLE),("S4","Layout + Proof","SceneLM + SceneProof",GREEN)]
    for i,(code,title,detail,color) in enumerate(stages):
        add_stage(s,xs[i],2.05,2.08,2.18,code,title,detail,color)
        if i < 4: add_arrow(s,xs[i]+2.10,3.14,xs[i+1]-0.05,3.14,MID,1.8)
    add_box(s, 0.72, 4.82, 11.88, 1.22, NAVY, NAVY)
    add_runs(s, [("输出  ",TEAL,True),("placement.json  ·  render.png  ·  evaluation.json  ·  sceneproof-result.zip",WHITE,False)], 1.08, 5.10, 11.1, 0.35, 17, MSO_ANCHOR.MIDDLE)
    add_text(s, "核心转变：S0–S3 负责“看懂并恢复”，S4 负责“让结果物理上更可信并给出证书”。", 1.08, 5.58, 11.0, 0.28, 12, RGBColor(207,221,246), True)
    add_text(s, "技术来源：Lumenarium 基于 Imaginarium；S2 DeepSearch 由 Calvin Gu 及其团队贡献。", 1.08, 6.22, 11.0, 0.26, 11.2, MID, False, PP_ALIGN.CENTER)

    # 5 Timeline
    s = new_slide(prs, "两个月迭代路径：每一次失败都变成下一层能力", "从基线复现，到论文级评测，再到可用服务", 5, total)
    y = 3.25
    line = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), Inches(y), Inches(12.2), Inches(y)); set_line(line, NAVY2, 3)
    milestones = [(1.05,"第 1–2 周","复现与统一评测","Paper30 / 8000px+"),(3.35,"第 3 周","V1 → V3","支撑树与堆叠感知"),(5.65,"第 4–5 周","SceneLM","关系范围优化"),(7.95,"第 6 周","SceneProof","证书与回滚"),(10.25,"第 7–8 周","产品化","双 A10 Web/API")]
    colors=[BLUE,ORANGE,TEAL,GREEN,PURPLE]
    for i,(x,when,title,detail) in enumerate(milestones):
        dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y-0.17),Inches(0.34),Inches(0.34));set_fill(dot,colors[i]);dot.line.fill.background()
        add_text(s,when,x-0.28,1.62,1.65,0.28,12,colors[i],True,PP_ALIGN.CENTER)
        add_text(s,title,x-0.44,2.03,1.98,0.36,16,NAVY,True,PP_ALIGN.CENTER)
        add_text(s,detail,x-0.58,3.72,2.30,0.50,11,MID,False,PP_ALIGN.CENTER)
    add_box(s, 1.03, 5.05, 11.15, 0.88, PALE_TEAL, TEAL)
    add_text(s, "主线始终不变：用显式关系和可验证几何，替代“看起来差不多”的隐式启发式。", 1.32, 5.30, 10.55, 0.32, 18, DARK, True, PP_ALIGN.CENTER)

    # 6 Contribution one
    s = new_slide(prs, "创新一：把“支撑”写进场景表示", "V1 → V3：从独立物体姿态，升级为结构化 support tree", 6, total)
    add_box(s,0.62,1.42,5.78,4.95,WHITE,GRID)
    add_text(s,"过去：每个物体各自优化",0.94,1.72,4.8,0.36,20,NAVY,True)
    add_bullets(s,[("二维相似：","能找到像的资产，但不知道“放在哪一层”。"),("孤立姿态：","床、枕头、床头柜之间缺少可执行结构。"),("失败模式：","父对象缺失时硬索引，或把错误关系传到 S4。")],0.94,2.32,4.92,2.15,15,bullet_color=RED)
    add_text(s,"结果：恢复率可高，但物理含义不完整。",0.94,5.42,4.9,0.36,15,RED,True)
    add_box(s,6.67,1.42,6.03,4.95,PALE_TEAL,TEAL)
    add_text(s,"现在：关系决定自由度与检查方式",7.02,1.72,5.2,0.36,20,TEAL_DARK,True)
    add_text(s,"floor_0",7.15,5.28,1.15,0.34,14,NAVY,True,PP_ALIGN.CENTER)
    add_text(s,"double_bed_0",8.48,3.97,1.72,0.34,14,NAVY,True,PP_ALIGN.CENTER)
    add_text(s,"pillow_0",7.10,2.62,1.30,0.34,13,DARK,True,PP_ALIGN.CENTER)
    add_text(s,"pillow_1",9.64,2.62,1.30,0.34,13,DARK,True,PP_ALIGN.CENTER)
    add_text(s,"side_table_0",10.85,4.15,1.55,0.34,13,DARK,True,PP_ALIGN.CENTER)
    add_arrow(s,7.73,5.15,8.93,4.34,TEAL); add_arrow(s,9.03,3.90,7.79,3.02,TEAL); add_arrow(s,9.55,3.90,10.25,3.02,TEAL); add_arrow(s,8.05,5.22,11.50,4.56,TEAL)
    add_text(s,"父子支撑树 + floor / wall / ceiling 分流 + stack-aware 序列",7.05,5.83,5.18,0.30,12,TEAL_DARK,True,PP_ALIGN.CENTER)

    # 7 Baseline lesson
    s = new_slide(prs, "V3 证明“结构化恢复”有效，但也暴露了代价", "Recovery 最强 ≠ 最终系统最好：还要同时考虑速度与物理质量", 7, total)
    data=[["版本","Primary recovery","Primary parent","Physical macro","S4 / 全链路"],["V1","89.49%","89.32%","52.98%","Legacy SA"],["V3","91.40%","87.80%","52.14%","约 24 min/scene"],["V4 DeepSearch","88.22%","80.14%","54.58%","检索升级"],["V5-fast","88.22%","80.14%","62.10%","S4 192.9s"]]
    add_table(s,data,0.64,1.46,12.02,2.48,[2.15,1.65,1.60,1.65,2.15],12.2,{2:PALE_ORANGE,4:PALE_TEAL})
    add_box(s,0.66,4.34,5.80,1.80,PALE_ORANGE,ORANGE)
    add_text(s,"V3 的价值",0.94,4.64,2.0,0.30,17,ORANGE,True)
    add_text(s,"支撑感知把 Primary recovery 推到 91.40%，说明结构信息确实帮助“找回物体”。",0.94,5.10,5.05,0.66,14,DARK)
    add_box(s,6.78,4.34,5.88,1.80,PALE_TEAL,TEAL)
    add_text(s,"为什么继续做 V5",7.06,4.64,2.5,0.30,17,TEAL_DARK,True)
    add_text(s,"SA-5000 太慢，且高 recovery 并不自动带来低碰撞与稳定支撑。需要新的 S4 范式。",7.06,5.10,5.05,0.66,14,DARK)

    # 8 SceneLM
    s = new_slide(prs, "创新二（上）：SceneLM 让优化从“盲搜”变成“定向修改”", "Relation Program 把语言模型的推理范围限制在可审计的关系与自由度内", 8, total)
    add_box(s,0.62,1.42,3.17,4.90,WHITE,GRID)
    add_text(s,"1  编译关系",0.92,1.72,2.4,0.35,19,NAVY,True)
    add_text(s,"SUPPORT\nCOLLISION_EXCLUSION\nPLANE_ATTACH\nSEMANTIC",0.95,2.34,2.45,1.55,16,BLUE,True,font=MONO)
    add_text(s,"把自然语言场景图编译为显式约束程序。",0.94,4.35,2.42,0.82,13,MID)
    add_box(s,4.08,1.42,4.42,4.90,PALE_BLUE,BLUE)
    add_text(s,"2  LM 只提议必要改动",4.40,1.72,3.6,0.35,19,BLUE,True)
    add_text(s,"Δx = {Δt, ΔR}Ω",4.52,2.50,3.55,0.54,25,NAVY,True,PP_ALIGN.CENTER,font="Cambria Math")
    add_text(s,"Ω = 被违反关系涉及的对象与自由度",4.54,3.18,3.55,0.42,14,DARK,True,PP_ALIGN.CENTER)
    add_bullets(s,[("局部：","不重排整间房。"),("稀疏：","不触碰无关对象。"),("可解释：","每个 proposal 对应明确 residual。")],4.50,3.90,3.42,1.48,13,bullet_color=BLUE)
    add_box(s,8.80,1.42,3.90,4.90,PALE_TEAL,TEAL)
    add_text(s,"3  交给证书裁决",9.12,1.72,3.0,0.35,19,TEAL_DARK,True)
    add_text(s,"proposal ≠ commit",9.15,2.53,3.16,0.48,23,NAVY,True,PP_ALIGN.CENTER,font=MONO)
    add_text(s,"模型负责提出候选；\n几何与指标负责决定是否采用。",9.18,3.38,3.05,0.92,15,DARK,True,PP_ALIGN.CENTER)
    add_text(s,"这正是速度与可靠性可以同时提升的原因。",9.18,4.78,3.03,0.65,13,TEAL_DARK,True,PP_ALIGN.CENTER)
    add_arrow(s,3.80,3.76,4.05,3.76,MID); add_arrow(s,8.52,3.76,8.77,3.76,MID)

    # 9 SceneProof
    s = new_slide(prs, "创新二（下）：SceneProof 把每次修改变成“可证明事务”", "局部通过、全局非劣，才允许从 incumbent 提交为新场景", 9, total)
    steps=[("候选","LM / rule proposal",BLUE),("重建","true mesh / voxel",PURPLE),("局部门","接触、COM、碰撞",ORANGE),("全局门","family non-regression",TEAL),("提交/恢复","commit or incumbent",GREEN)]
    xs=[0.55,3.03,5.51,7.99,10.47]
    for i,(a,b,c) in enumerate(steps):
        add_box(s,xs[i],1.55,2.12,1.25,WHITE,c)
        add_text(s,a,xs[i]+0.18,1.77,1.76,0.31,16,c,True,PP_ALIGN.CENTER)
        add_text(s,b,xs[i]+0.15,2.20,1.82,0.26,10.2,MID,False,PP_ALIGN.CENTER)
        if i<4:add_arrow(s,xs[i]+2.14,2.17,xs[i+1]-0.04,2.17,MID,1.8)
    add_box(s,0.68,3.26,5.90,2.62,WHITE,GRID)
    add_text(s,"局部证书",0.98,3.56,2.0,0.34,18,NAVY,True)
    add_bullets(s,[("支撑：","declared parent / sibling contact + COM margin"),("碰撞：","exact-mesh overlap 不增加；穿透受限"),("姿态：","冻结无关 SO(3) / 高度 / 切向自由度"),("可见性：","S1 mask reprojection 不明显恶化")],0.98,4.03,5.10,1.55,13.2)
    add_box(s,6.84,3.26,5.82,2.62,PALE_TEAL,TEAL)
    add_text(s,"提交准则",7.15,3.56,2.0,0.34,18,TEAL_DARK,True)
    add_text(s,"Accept(Δ) ⇔ LocalGates(Δ) ∧ ΔFamily ≥ −ε",7.18,4.12,5.08,0.48,20,NAVY,True,PP_ALIGN.CENTER,font="Cambria Math")
    add_text(s,"任何失败 → 恢复当前 incumbent；\n证据不足 → 标记 unresolved，而不是伪装成功。",7.22,4.88,5.02,0.66,14,DARK,True,PP_ALIGN.CENTER)

    # 10 Failure-driven iteration
    s = new_slide(prs, "一次枕头问题，推动出完整的“支撑证明”", "从简单掉落，到 constrained projection、first-contact 与回滚", 10, total)
    add_picture_cover(s, before, 0.60, 1.42, 5.76, 4.75)
    add_picture_cover(s, after, 6.72, 1.42, 5.96, 4.75)
    add_box(s,0.85,5.35,4.95,0.60,PALE_RED,RED); add_text(s,"问题：悬空 / 超出支撑 / 错误旋转",1.04,5.52,4.60,0.26,14,RED,True,PP_ALIGN.CENTER)
    add_box(s,7.01,5.35,5.35,0.60,PALE_TEAL,TEAL); add_text(s,"策略：最小切向投影 → Z-only first contact → 回滚",7.20,5.52,4.98,0.26,13.3,TEAL_DARK,True,PP_ALIGN.CENTER)
    add_text(s,"工程经验：物理修复不能靠“多掉一点”；必须知道真实支撑面、保留哪些自由度，以及何时停止。",1.15,6.42,11.0,0.34,15,DARK,True,PP_ALIGN.CENTER)

    # 11 Quality
    s = new_slide(prs, "结果一：V5-fast 在保持上游工作点时，显著改善物理质量", "Paper30；Primary 只统计 S1 可见掩码 ≥ 8,000 px 的对象；GT 只用于评测", 11, total)
    data=[["版本","Primary recovery","Primary parent","Physical macro","说明"],["Imaginarium V1","89.49%","89.32%","52.98%","原始基线"],["Lumenarium V3","91.40%","87.80%","52.14%","支撑感知恢复"],["V4 DeepSearch","88.22%","80.14%","54.58%","V5 的上游输入"],["V5-fast / Fix61","88.22%","80.14%","62.10%","论文主版本"]]
    add_table(s,data,0.68,1.48,11.98,2.60,[2.4,1.8,1.7,1.65,2.0],12.5,{4:PALE_TEAL})
    add_metric(s,0.70,4.45,3.65,1.48,"+7.52 pp","Physical macro","相对 V4 DeepSearch；主要来自 SceneLM + SceneProof",GREEN)
    add_box(s,4.67,4.45,3.82,1.48,PALE_BLUE,BLUE)
    add_text(s,"关键归因",4.96,4.73,2.2,0.30,17,BLUE,True)
    add_text(s,"V4 与 V5 的 recovery / parent 相同，因此物理增益不是“上游多找回了物体”。",4.96,5.18,3.13,0.63,12.2,DARK)
    add_box(s,8.80,4.45,3.84,1.48,PALE_ORANGE,ORANGE)
    add_text(s,"论文边界",9.08,4.73,2.2,0.30,17,ORANGE,True)
    add_text(s,"V5-medium 属于 presentation-only visual cleanup，不混入主表。",9.08,5.18,3.12,0.63,12.2,DARK)

    # 12 Speed
    s = new_slide(prs, "结果二：S4 已提速 3.51×，端到端瓶颈转移到 S1", "全部数字来自 Paper30 冷启动与同输入 S4 基准", 12, total)
    add_picture_cover(s,speed_chart,0.60,1.42,7.20,3.12,border=WHITE,border_width=0)
    add_box(s,8.12,1.44,4.54,3.08,WHITE,GRID)
    add_text(s,"S4 同输入测速",8.43,1.76,3.6,0.33,19,NAVY,True)
    add_text(s,"677.770 s",8.44,2.42,1.48,0.36,21,MID,True)
    add_text(s,"Legacy SA-5000",9.94,2.48,2.12,0.26,11,MID)
    add_text(s,"↓ 3.513×",8.45,3.00,2.0,0.34,20,TEAL,True)
    add_text(s,"192.930 s",8.44,3.50,1.48,0.36,21,GREEN,True)
    add_text(s,"SceneLM + Fix61",9.94,3.56,2.12,0.26,11,MID)
    add_metric(s,0.72,5.05,3.65,1.18,"829.879 s","V5-fast 全链路","13.83 min/scene；S0–S4 实测均值",TEAL)
    add_metric(s,4.60,5.05,3.65,1.18,"2.680 h","S0–S3 双 A10","Paper30 实测墙钟时间",BLUE)
    add_metric(s,8.48,5.05,3.65,1.18,"250–320 s","S1 目标区间","若 Gemini 稳定 8 并发；容量规划估算",ORANGE)

    # 13 Profiles
    s = new_slide(prs, "同一套底座，面向三种使用场景", "Fast 保论文可比性；Medium 优先展示安全；Best 用多冷启换上限", 13, total)
    cards=[(0.62,"V5-fast","Fix61 quantitative",TEAL,"论文定量 / 快速预览",["单次冷启动","证书化 SceneLM 结果","不做展示性删除"]),(4.52,"V5-medium","Fix61 + visual-safe",ORANGE,"技术美术 / 演示交付",["复用冻结上游","保守落地或隐藏少量叶子重复物","presentation-only"]),(8.42,"V5-best","3 seeds + selector",PURPLE,"最终精选 / 最高质量",["双 A10 并行三次冷启","GT-free selector","证书 → unresolved → 碰撞 → physical → coverage"])]
    for x,name,sub,c,use,bul in cards:
        add_box(s,x,1.48,3.62,4.78,WHITE,c,width=1.7)
        add_text(s,name,x+0.28,1.77,3.0,0.42,23,c,True)
        add_text(s,sub,x+0.30,2.25,2.9,0.26,11,MID,False,font=MONO)
        add_box(s,x+0.28,2.75,3.02,0.52,RGBColor(247,249,252),c)
        add_text(s,use,x+0.42,2.90,2.72,0.22,12,DARK,True,PP_ALIGN.CENTER)
        add_bullets(s,bul,x+0.32,3.62,2.95,1.70,12.5,bullet_color=c)
    add_text(s,"选择原则：定量系统与展示策略分开命名、分开评测、保留完整 provenance。",1.18,6.53,10.9,0.30,15,NAVY,True,PP_ALIGN.CENTER)

    # 14 Productization
    s = new_slide(prs, "研究不止停在脚本：已封装为双 A10 在线服务", "https://embedding.lightart.qq.com/ · 上传 1024×1024 PNG/JPEG", 14, total)
    add_box(s,0.62,1.44,3.00,4.78,WHITE,GRID)
    add_text(s,"用户体验",0.94,1.75,2.1,0.34,19,NAVY,True)
    add_bullets(s,["上传图片","选择 Fast / Medium / Best","实时查看 S0–S4 进度","下载 ZIP 与证书"],0.94,2.32,2.2,2.52,14)
    add_box(s,3.93,1.44,4.52,4.78,PALE_BLUE,BLUE)
    add_text(s,"服务架构",4.25,1.75,2.2,0.34,19,BLUE,True)
    add_text(s,"Web / REST API",5.02,2.37,2.34,0.42,17,NAVY,True,PP_ALIGN.CENTER)
    add_arrow(s,6.18,2.87,6.18,3.34,BLUE)
    add_text(s,"原子 Job Store + Cache",4.72,3.42,2.94,0.42,16,NAVY,True,PP_ALIGN.CENTER)
    add_arrow(s,5.42,3.98,4.70,4.49,BLUE); add_arrow(s,6.92,3.98,7.62,4.49,BLUE)
    add_text(s,"GPU 0 worker",4.13,4.61,1.88,0.38,14,DARK,True,PP_ALIGN.CENTER)
    add_text(s,"GPU 1 worker",6.37,4.61,1.88,0.38,14,DARK,True,PP_ALIGN.CENTER)
    add_text(s,"DeepSearch 最多 8 路聚合并发",4.53,5.38,3.35,0.30,12,BLUE,True,PP_ALIGN.CENTER)
    add_box(s,8.76,1.44,3.94,4.78,WHITE,GRID)
    add_text(s,"为真实运行补齐的可靠性",9.07,1.75,3.08,0.34,18,NAVY,True)
    add_bullets(s,[("缓存：","相同图跨 profile 复用冻结 S0–S3/Fix61。"),("调度：","原子 claim，避免双 GPU 重复领取。"),("恢复：","死亡 claim 回收、单场景失败不终止队列。"),("重试：","HTTP 流中断、S4 输出缺失有限重试。"),("可追溯：","release、seed、耗时、证书写入结果包。")],9.06,2.30,3.08,2.98,12.3)

    # 15 Demo
    s = new_slide(prs, "真实输出：从 A10 服务直接生成，而非人工搭建", "以下图片已放入工蜂 README；完整演示见 Bilibili BV1tpbD6hERB", 15, total)
    add_picture_cover(s,office,0.60,1.42,6.05,4.72)
    add_picture_cover(s,living,6.72,1.42,5.96,4.72)
    add_text(s,"办公场景",0.62,6.25,6.00,0.28,13,NAVY,True,PP_ALIGN.CENTER)
    add_text(s,"客厅场景",6.74,6.25,5.92,0.28,13,NAVY,True,PP_ALIGN.CENTER)
    add_box(s,3.90,6.57,5.57,0.42,NAVY,NAVY)
    add_text(s,"在线服务：embedding.lightart.qq.com",4.06,6.67,5.24,0.20,11,WHITE,True,PP_ALIGN.CENTER)

    # 16 What we learned
    s = new_slide(prs, "最重要的技术判断：把“可见效果”和“论文证据”分开", "两个月里真正困难的，不是写一个 loss，而是做正确归因与停止错误提交", 16, total)
    add_box(s,0.65,1.45,3.76,4.82,WHITE,RED)
    add_text(s,"1  不静默成功",0.96,1.77,2.5,0.34,19,RED,True)
    add_text(s,"修复失败时标记 unresolved；\n不能因为渲染完成就宣称物理通过。",0.97,2.43,2.94,0.95,15,DARK,True)
    add_text(s,"领导视角：风险可见，才可交付。",0.98,5.42,2.92,0.36,13,RED,True)
    add_box(s,4.78,1.45,3.76,4.82,WHITE,ORANGE)
    add_text(s,"2  不错误归因",5.09,1.77,2.5,0.34,19,ORANGE,True)
    add_text(s,"V3 → V4 的 pose 降幅出现在 DeepSearch 上游；V4 → V5 基本不变。\n因此不能归因给 SceneLM。",5.10,2.43,2.94,1.25,14.5,DARK,True)
    add_text(s,"论文视角：同输入、冻结上游、分层对比。",5.10,5.42,2.94,0.36,13,ORANGE,True)
    add_box(s,8.91,1.45,3.76,4.82,WHITE,TEAL)
    add_text(s,"3  不混淆产品模式",9.22,1.77,2.7,0.34,19,TEAL_DARK,True)
    add_text(s,"Fast 用于定量；Medium 可做保守视觉清理；Best 用无 GT selector 选冷启。",9.23,2.43,2.94,1.15,14.5,DARK,True)
    add_text(s,"产品视角：质量、速度与证据强度可选择。",9.23,5.42,2.94,0.36,13,TEAL_DARK,True)

    # 17 Next steps
    s = new_slide(prs, "下一步：从“可用”走向“规模化与更高上限”", "优先攻端到端瓶颈，同时保持证书边界不后退", 17, total)
    add_box(s,0.65,1.48,7.42,4.72,WHITE,GRID)
    add_text(s,"近期（可直接落地）",0.98,1.79,3.2,0.34,19,NAVY,True)
    add_bullets(s,[("S1 并发：","解除全局 API lock，验证 Gemini 8 并发，把 S1 由 443s 推向 250–320s。"),("DeepSearch 校准：","做 seed-locked S2-only ablation，定位资产 frame / scale / parent 偏差。"),("V5-best：","完善三冷启无 GT selector，报告额外算力与收益曲线。"),("可观测性：","把每阶段耗时、GPU-hours、unresolved 和失败原因统一进 dashboard。")],0.98,2.35,6.50,2.90,14.2)
    add_box(s,8.38,1.48,4.29,4.72,PALE_TEAL,TEAL)
    add_text(s,"中期研究方向",8.71,1.79,3.2,0.34,19,TEAL_DARK,True)
    add_bullets(s,[("Pose：","Flux fine-tuning / retrieval-frame recalibration。"),("Proof：","更快的 dominant support component 与 exact-mesh narrow phase。"),("Tooling：","UE / DCC 导入、人工确认点、可编辑场景回写。")],8.71,2.36,3.32,2.42,14)
    add_box(s,8.71,5.20,3.58,0.58,NAVY,NAVY)
    add_text(s,"原则：先减少最大瓶颈，再增加模型规模。",8.88,5.37,3.26,0.23,11.5,WHITE,True,PP_ALIGN.CENTER)

    # 18 Closing
    s = prs.slides.add_slide(prs.slide_layouts[6]); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    band=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,Inches(0.14));set_fill(band,TEAL);band.line.fill.background()
    add_text(s,"Lumenarium",0.88,0.72,5.0,0.48,28,TEAL,True)
    add_text(s,"我的暑期实习，不只是把一个指标做高",0.88,1.55,10.8,0.55,29,WHITE,True)
    add_text(s,"而是把“单图 3D 场景生成”推进成一套更快、可验证、可部署的系统。",0.90,2.35,11.35,0.52,24,WHITE,True)
    add_box(s,0.90,3.42,11.56,1.43,NAVY2,TEAL)
    add_runs(s,[("结构化恢复  ",TEAL,True),("→ 关系范围优化  ",WHITE,False),("→ Proof-carrying commit  ",WHITE,False),("→ 双 A10 服务",WHITE,False)],1.25,3.81,10.85,0.52,19,MSO_ANCHOR.MIDDLE,PP_ALIGN.CENTER)
    add_text(s,"谢谢 · Questions",0.90,5.70,5.0,0.46,24,WHITE,True)
    add_text(s,"Hansen Zhu  ·  Mentor: Calvin Gu",0.92,6.32,5.4,0.28,13,RGBColor(205,221,247))
    add_text(s,"Demo: https://www.bilibili.com/video/BV1tpbD6hERB/",0.92,6.73,7.0,0.23,10.5,RGBColor(205,221,247))
    add_text(s,"Foundation: Imaginarium (SIGGRAPH Asia 2025) · DeepSearch: Calvin Gu and team",6.55,6.72,5.82,0.24,9.2,RGBColor(205,221,247),False,PP_ALIGN.RIGHT)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
