from __future__ import annotations

from pathlib import Path
from math import ceil

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\hansenzhu\Desktop\ppt_images\Lumenarium")
TEMPLATE = ROOT.parent / "Imaginarium_v1_v3_Professional.pptx"
RENDERS = ROOT / "sceneproof_paper30_comparison_renders_fix25"
OUT = ROOT / "SceneProof_Weekly_Report_2026-08-04.pptx"
ASSETS = ROOT / "weekly_report_assets"

NAVY = RGBColor(18, 38, 95)
NAVY_2 = RGBColor(27, 61, 115)
TEAL = RGBColor(0, 196, 157)
TEAL_DARK = RGBColor(0, 142, 118)
PALE_TEAL = RGBColor(228, 247, 241)
LIGHT = RGBColor(244, 246, 251)
MID = RGBColor(105, 114, 132)
DARK = RGBColor(29, 45, 88)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(229, 66, 66)
PALE_RED = RGBColor(255, 239, 239)
ORANGE = RGBColor(244, 153, 22)
GREEN = RGBColor(32, 157, 91)
BLUE = RGBColor(42, 99, 190)
GRID = RGBColor(214, 220, 231)
FONT = "Microsoft YaHei"
MONO = "Consolas"


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst  # noqa: SLF001 - intentional template reuse
    for slide_id in list(slide_ids):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_ids.remove(slide_id)


def set_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor, width: float = 1.0, transparency: int = 0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    shape.line.transparency = transparency


def add_text(slide, text, x, y, w, h, size=20, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, valign=MSO_ANCHOR.TOP,
             margin=0.02, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_runs(slide, runs, x, y, w, h, size=19, color=DARK,
             valign=MSO_ANCHOR.TOP, line_spacing=1.08):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.line_spacing = line_spacing
    for item in runs:
        text, item_color, bold = item[:3]
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = item_color or color
    return box


def add_box(slide, x, y, w, h, fill=WHITE, line=GRID, radius=True, width=1.2):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, width)
    return shape


def add_title(slide, title, subtitle, page, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(0.43), Inches(0.16), Inches(0.63))
    set_fill(bar, TEAL)
    bar.line.fill.background()
    add_text(slide, title, 0.9, 0.40, 10.8, 0.48, 26, NAVY, True)
    add_text(slide, subtitle, 0.9, 0.88, 10.7, 0.28, 12, MID)
    add_text(slide, f"{page:02d} / {total:02d}", 11.78, 0.48, 0.95, 0.28, 11, MID, True, PP_ALIGN.RIGHT)


def add_footer(slide, text="SceneProof · Weekly Research Update · 2026-08-04"):
    add_text(slide, text, 0.65, 7.23, 12.05, 0.18, 8, MID)


def add_bullets(slide, items, x, y, w, h, size=17, color=DARK,
                bullet_color=TEAL, spacing=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        if isinstance(item, tuple):
            head, body = item
        else:
            head, body = "", item
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.line_spacing = 1.08
        r = p.add_run()
        r.text = "▌ "
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = bullet_color
        if head:
            r = p.add_run()
            r.text = head
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = color
        r = p.add_run()
        r.text = body
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_table(slide, data, x, y, w, h, col_widths=None, font_size=12,
              header_fill=NAVY_2, highlight_rows=None, header_size=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(Inches(w) * cw / total)
    highlight_rows = highlight_rows or {}
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                text_color = WHITE
                bold = True
            else:
                fill = highlight_rows.get(r, WHITE if r % 2 else LIGHT)
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill
                text_color = DARK
                bold = c == 0
            cell.text = str(data[r][c])
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            p.vertical_anchor = MSO_ANCHOR.MIDDLE
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(header_size or font_size) if r == 0 else Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = text_color
    return table_shape


def add_picture_cover(slide, path: Path, x, y, w, h, border=GRID, border_width=1.0):
    with Image.open(path) as im:
        iw, ih = im.size
    target = w / h
    source = iw / ih
    crop_l = crop_r = crop_t = crop_b = 0.0
    if source > target:
        visible = target / source
        crop_l = crop_r = (1 - visible) / 2
    else:
        visible = source / target
        crop_t = crop_b = (1 - visible) / 2
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    pic.crop_left = crop_l
    pic.crop_right = crop_r
    pic.crop_top = crop_t
    pic.crop_bottom = crop_b
    set_line(pic, border, border_width)
    return pic


def make_support_limitation_crop() -> Path:
    src = RENDERS / "bedroom_01" / "03_certified_final.png"
    out = ASSETS / "bedroom_01_support_limitation.png"
    with Image.open(src).convert("RGB") as im:
        # Crop the bed / pillows while retaining enough support context.
        crop = im.crop((500, 390, 930, 790)).resize((860, 800), Image.Resampling.LANCZOS)
        draw = ImageDraw.Draw(crop)
        draw.ellipse((560, 155, 825, 345), outline=(230, 60, 60), width=10)
        draw.line((760, 125, 690, 205), fill=(230, 60, 60), width=12)
        draw.polygon([(690, 205), (710, 170), (730, 207)], fill=(230, 60, 60))
        crop.save(out, quality=95)
    return out


def add_metric_card(slide, x, y, w, h, value, label, accent=TEAL, note=None):
    add_box(slide, x, y, w, h, WHITE, GRID)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    set_fill(band, accent)
    band.line.fill.background()
    add_text(slide, value, x + 0.25, y + 0.18, w - 0.38, 0.46, 27, accent, True)
    add_text(slide, label, x + 0.25, y + 0.67, w - 0.38, 0.30, 12, DARK, True)
    if note:
        add_text(slide, note, x + 0.25, y + 1.02, w - 0.38, h - 1.12, 9.5, MID)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = NAVY
    # Layered geometry recreates the prior report's title language.
    poly = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(8.0), Inches(-0.2), Inches(5.8), Inches(7.9))
    set_fill(poly, NAVY_2, 25); poly.line.fill.background()
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14))
    set_fill(top, TEAL); top.line.fill.background()
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.36), prs.slide_width, Inches(0.14))
    set_fill(bot, TEAL); bot.line.fill.background()
    add_text(slide, "WEEKLY RESEARCH UPDATE · 2026-08-04", 0.9, 0.72, 8.9, 0.35, 15, TEAL, True)
    add_text(slide, "SceneProof", 0.9, 1.55, 7.8, 0.72, 43, WHITE, True)
    add_text(slide, "关系程序驱动的可认证场景优化", 0.9, 2.30, 9.8, 0.64, 29, WHITE, True)
    add_text(slide, "Full-SO(3) · Topology-aware Schur · Component-wise Certificate", 0.93, 3.04, 10.5, 0.38, 17, RGBColor(199, 220, 251))
    add_box(slide, 0.9, 4.05, 5.75, 1.08, NAVY_2, TEAL)
    add_runs(slide, [
        ("Paper30 结论  ", TEAL, True),
        ("S4 3.10×；S2 8-way mean 1.274×；fix25 visibility certificate 待修", WHITE, False),
    ], 1.13, 4.28, 5.3, 0.62, 16, WHITE, MSO_ANCHOR.MIDDLE)
    add_text(slide, "共同输入：frozen v4-deepsearch S0–S3  |  S4-only paired evaluation", 0.93, 5.43, 8.8, 0.30, 13, RGBColor(207, 221, 246))
    add_text(slide, "本周：从通用优化器推进到 proof-carrying scene backend", 0.93, 6.23, 8.8, 0.34, 15, WHITE, True)


def build():
    ASSETS.mkdir(exist_ok=True)
    support_crop = make_support_limitation_crop()
    prs = Presentation(str(TEMPLATE))
    remove_all_slides(prs)
    prs.core_properties.title = "SceneProof Weekly Research Update"
    prs.core_properties.subject = "Paper30 full-SO(3) guarded Schur and component certificate"
    prs.core_properties.author = "Lumenarium"
    total = 12
    title_slide(prs)

    # 2 Executive summary
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "本周结论：速度与物理质量同时成立", "Paper30 · 30/30 scenes · 10,000× paired bootstrap · fixed camera renders", 2, total)
    add_metric_card(slide, 0.60, 1.42, 2.75, 1.52, "3.10×", "S4 speedup vs Legacy SA5000", TEAL, "677.77s → 218.80s / scene")
    add_metric_card(slide, 3.50, 1.42, 2.75, 1.52, "1.274×", "S2 mean speedup · 8-way", ORANGE, "83.97s → 65.90s；median 反而 81.36s → 92.19s")
    add_metric_card(slide, 6.40, 1.42, 2.75, 1.52, "+0.286 pp", "Physical macro vs smooth control", GREEN, "95% CI [+0.070, +0.538] pp")
    add_metric_card(slide, 9.30, 1.42, 3.38, 1.52, "Non-inferior", "Rotation / Translation AUC", BLUE, "collision +0.982 pp · pose 非劣")
    add_box(slide, 0.60, 3.28, 7.32, 3.43, WHITE, GRID)
    add_text(slide, "这次真正成功的是什么", 0.88, 3.55, 4.5, 0.36, 20, NAVY, True)
    add_bullets(slide, [
        ("结构性算法：", "Relation Program → sparse factor ownership → guarded Schur。"),
        ("安全性：", "每个物理分量单独门控；失败组件回滚，不以 macro 掩盖局部退化。"),
        ("公平性：", "v5 与 control 共用 frozen v4-deepsearch S3；GT 仅用于事后评估。"),
        ("Recovery：", "资产/父子图未改，因此 matched、recovery、parent 指标按设计保持不变。"),
    ], 0.90, 4.02, 6.78, 2.40, 15.5)
    add_box(slide, 8.18, 3.28, 4.50, 3.43, PALE_RED, RED)
    add_text(slide, "仍未解决", 8.48, 3.55, 2.2, 0.34, 20, RED, True)
    add_text(slide, "接触 ≠ 稳定支撑", 8.48, 4.08, 3.65, 0.36, 20, NAVY, True)
    add_text(slide, "枕头虽与床相交/接触，但投影重心可能落在支撑多边形之外；当前静态 overlap 不能证明抗扰稳定性。", 8.48, 4.58, 3.78, 1.05, 14, DARK)
    add_text(slide, "下一步：COM margin + support patch + perturbation survival", 8.48, 5.82, 3.70, 0.55, 14, RED, True)
    add_footer(slide)

    # 3 Data lineage / upstream table
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "数据谱系：v5 使用 v4-deepsearch 的 S3", "不要把上游 S1–S3 差异与本次 S4 solver 增益混为一谈", 3, total)
    data = [
        ["上游版本", "Scenes", "Matched/GT", "Primary recovery", "Primary parent", "Rot AUC@60", "Trans AUC@0.5m"],
        ["v1", "30/30", "307/363", "0.8949", "0.8932", "0.4813", "0.2373"],
        ["v3", "30/30", "旧 frozen run", "0.9108", "0.9394", "0.4594", "0.2063"],
        ["v4-deepsearch", "30/30", "302/363", "0.8822", "0.8014", "0.3134", "0.1219"],
    ]
    add_table(slide, data, 0.62, 1.44, 12.05, 2.22, [2.0,1.1,1.45,1.65,1.55,1.4,1.55], 12.5,
              highlight_rows={3: PALE_TEAL})
    add_box(slide, 0.62, 3.97, 7.28, 2.67, WHITE, GRID)
    add_text(slide, "本次 Paper30 固定协议", 0.92, 4.22, 4.3, 0.34, 20, NAVY, True)
    add_bullets(slide, [
        ("共同输入：", "v4-deepsearch 的 S0–S3、资产缓存、frozen S3 geometry。"),
        ("Legacy 对照：", "同一输入上的 SA5000；仅用于速度与旧 S4 对照。"),
        ("v5 质量对照：", "v5_sceneproof_smooth_control_paper30_fix25。"),
        ("最终方法：", "v5_sceneproof_postsim_component_certified_paper30_fix25。"),
    ], 0.94, 4.70, 6.65, 1.65, 14.5)
    add_box(slide, 8.18, 3.97, 4.49, 2.67, PALE_TEAL, TEAL)
    add_text(slide, "Recovery 为什么不变？", 8.48, 4.22, 3.6, 0.34, 20, TEAL_DARK, True)
    add_text(slide, "v5 只优化 S4 continuous pose；不替换 asset，不改 parent，不增删对象。", 8.48, 4.78, 3.70, 0.72, 15, DARK)
    add_text(slide, "因此 302/363 matched、0.8822 primary recovery、0.8014 primary parent 由共同 S3 决定。", 8.48, 5.53, 3.70, 0.74, 14, NAVY, True)
    add_text(slide, "注：当前本地冻结汇总未保存 v1/v3 的 secondary split，周报不补造数值。", 0.68, 6.78, 11.9, 0.24, 9.5, MID)
    add_footer(slide)

    # 4 Inspiration and manifold
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "算法灵感：有效场景位于关系条件化的分层流形", "不是把所有物体当作彼此独立的 6-DoF 块", 4, total)
    add_box(slide, 0.62, 1.43, 5.78, 4.93, PALE_RED, RED)
    add_text(slide, "Ambient-space optimizer", 0.95, 1.72, 4.5, 0.35, 22, RED, True)
    add_text(slide, "𝒳 = ∏ᵢ (SO(3) × ℝ³)", 1.08, 2.31, 4.80, 0.48, 27, NAVY, True, PP_ALIGN.CENTER, "Cambria Math")
    add_bullets(slide, [
        ("自由度过多：", "噪声会沿不受约束方向漂移。"),
        ("耦合未利用：", "桌面移动与桌上物体被重复求解。"),
        ("非光滑：", "碰撞/contact mode 激活时，单一 LM/Adam 假设失效。"),
        ("输出无证明：", "较低总 loss 不能说明每个物理分量都未退化。"),
    ], 0.95, 3.05, 4.98, 2.68, 15)
    add_box(slide, 6.66, 1.43, 6.02, 4.93, PALE_TEAL, TEAL)
    add_text(slide, "SceneProof view", 7.00, 1.72, 4.0, 0.35, 22, TEAL_DARK, True)
    add_text(slide, "𝓜 = ⋃σ 𝓜σ  ⊂ 𝒳", 7.17, 2.31, 4.95, 0.48, 27, NAVY, True, PP_ALIGN.CENTER, "Cambria Math")
    add_bullets(slide, [
        ("局部流形：", "固定 support/contact mode σ 时，关系定义可行切空间。"),
        ("分层结构：", "碰撞激活与关系释放对应 strata 之间的切换。"),
        ("自适应 chart：", "可靠关系采用低维 translation chart；不可靠块保持自由。"),
        ("Fail-closed：", "跨 strata 的失败候选回滚到 incumbent。"),
    ], 7.00, 3.05, 5.10, 2.68, 15)
    add_text(slide, "关键取舍：full SO(3) 保留表达能力；只压缩经过审计的 leaf translation。", 1.0, 6.55, 11.2, 0.36, 17, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 5 Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "SceneProof-Core：从关系图到 proof-carrying scene", "同一 Relation Program 同时生成 solver factor、ownership、probe 与 certificate", 5, total)
    stages = [
        ("Frozen S3", "asset / parent / scale\npose + geometry"),
        ("Program IR", "SUPPORT · PLANE\nALIGN · DISTANCE"),
        ("Sparse chart", "full SO(3) root\nrelation translation"),
        ("Guarded solve", "Schur + back-sub\nall-collision check"),
        ("Certificate", "component accept\nor exact rollback"),
    ]
    x0 = 0.63
    for i, (head, body) in enumerate(stages):
        x = x0 + i * 2.55
        fill = PALE_TEAL if i in (1, 2, 4) else WHITE
        line = TEAL if i in (1, 2, 4) else GRID
        add_box(slide, x, 1.67, 2.05, 1.55, fill, line)
        add_text(slide, head, x + 0.13, 1.88, 1.78, 0.30, 17, NAVY, True, PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.13, 2.29, 1.78, 0.60, 11.5, DARK, False, PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_text(slide, "→", x + 2.09, 2.11, 0.40, 0.40, 23, TEAL, True, PP_ALIGN.CENTER)
    add_box(slide, 0.63, 3.63, 12.04, 2.65, WHITE, GRID)
    add_text(slide, "Program contract", 0.93, 3.92, 2.4, 0.32, 20, NAVY, True)
    add_text(slide, "P_f  ↦  ( r_f, owners(f), 𝒯_f, probe_f, witness_f )", 3.10, 3.86, 8.75, 0.46, 24, NAVY, True, PP_ALIGN.CENTER, "Cambria Math")
    add_bullets(slide, [
        ("Residual：", "数值目标与可微/可评估语义；"),
        ("Ownership：", "Jacobian 只允许连接声明过的变量块；"),
        ("Tangent chart：", "决定哪些 translation 可安全消元；"),
        ("Probe + witness：", "决定候选能否提交，以及失败时释放哪一局部组件。"),
    ], 0.95, 4.54, 11.0, 1.35, 14.5)
    add_footer(slide)

    # 6 math coordinates
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "独立 full-SO(3) + 关系条件化平移", "避免 parent yaw 错误传播，同时利用 support topology", 6, total)
    add_box(slide, 0.62, 1.43, 5.92, 4.95, WHITE, GRID)
    add_text(slide, "Rotation 始终在 root system", 0.94, 1.74, 4.9, 0.36, 21, NAVY, True)
    add_text(slide, "Rᵢ(ωᵢ) = Rᵢ⁰ exp([ωᵢ]×)", 1.06, 2.42, 4.90, 0.50, 26, TEAL_DARK, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "每个 object 保留独立世界坐标 SO(3)。support child 不默认继承 parent yaw；rotation parameters 永不 Schur eliminate。", 1.02, 3.15, 4.98, 1.13, 15, DARK)
    add_box(slide, 1.02, 4.58, 5.10, 1.23, PALE_TEAL, TEAL)
    add_text(slide, "好处：表达能力 ≈ full 6-DoF；减少过去 yaw inheritance 对 pose AUC 的伤害。", 1.26, 4.86, 4.62, 0.64, 14, NAVY, True)
    add_box(slide, 6.78, 1.43, 5.90, 4.95, WHITE, GRID)
    add_text(slide, "Translation 根据关系切换 chart", 7.10, 1.74, 5.0, 0.36, 21, NAVY, True)
    add_text(slide, "tᵢ = tₚ₍ᵢ₎ + Bᵢuᵢ", 7.26, 2.42, 4.90, 0.50, 26, TEAL_DARK, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "兼容且线性化稳定的 support leaf 使用关系坐标 uᵢ；关系不可靠、collision separator 或非 leaf 仍保留完整自由变量。", 7.15, 3.15, 4.94, 1.13, 15, DARK)
    add_box(slide, 7.15, 4.58, 5.10, 1.23, PALE_TEAL, TEAL)
    add_text(slide, "好处：只去掉无效/重复自由度，而不是用错误先验把场景锁死。", 7.39, 4.86, 4.62, 0.64, 14, NAVY, True)
    add_footer(slide)

    # 7 Schur math
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "拓扑 Schur：先解 root，再恢复 leaf", "对当前线性化子问题是精确消元，不是近似删变量", 7, total)
    add_box(slide, 0.62, 1.43, 7.33, 4.95, WHITE, GRID)
    add_text(slide, "线性化与阻尼正规方程", 0.94, 1.72, 4.1, 0.34, 20, NAVY, True)
    add_text(slide, "r(x ⊞ δ) ≈ r(x) + Jδ", 1.15, 2.23, 5.90, 0.42, 23, TEAL_DARK, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "[ Hrr  Hrl ] [δr] = −[gr]\n[ Hlr  Hll ] [δl]    [gl]", 1.42, 2.90, 5.36, 1.10, 22, NAVY, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "S = Hrr − Hrl Hll⁻¹ Hlr", 1.15, 4.16, 5.90, 0.42, 22, TEAL_DARK, True, PP_ALIGN.CENTER, "Cambria Math")
    add_text(slide, "Sδr = −(gr − Hrl Hll⁻¹gl)\nδl = −Hll⁻¹(gl + Hlrδr)", 1.08, 4.77, 6.04, 0.84, 20, NAVY, True, PP_ALIGN.CENTER, "Cambria Math")
    add_box(slide, 8.22, 1.43, 4.46, 4.95, PALE_TEAL, TEAL)
    add_text(slide, "为什么更快", 8.52, 1.72, 2.8, 0.34, 20, TEAL_DARK, True)
    add_bullets(slide, [
        ("Topology sparsity：", "factor 只连接 relation neighborhood。"),
        ("Leaf elimination：", "避免反复求解弱观测的 child translation。"),
        ("Small root：", "所有 SO(3) 仍保留，但 translation system 更小、更稳定。"),
        ("Two-branch budget：", "短程 smooth + guarded candidate + certificate。"),
    ], 8.50, 2.28, 3.65, 2.95, 14.5)
    add_text(slide, "关键：仅消元 audit 通过的 leaf；否则 abstain。", 8.54, 5.55, 3.56, 0.50, 14, RED, True)
    add_footer(slide)

    # 8 Certificate
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "Component-wise certificate：macro 不能掩盖局部退化", "候选生成与候选提交被明确分离", 8, total)
    add_box(slide, 0.62, 1.45, 4.03, 4.95, WHITE, GRID)
    add_text(slide, "1 · Candidate", 0.94, 1.75, 2.6, 0.34, 20, NAVY, True)
    add_text(slide, "Schur / back-sub / poll", 1.02, 2.32, 3.22, 0.38, 18, TEAL_DARK, True, PP_ALIGN.CENTER)
    add_text(slide, "产生 x⁺，但不立即覆盖 incumbent x⁰。", 1.02, 2.92, 3.20, 0.70, 15, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, "↓", 2.15, 3.76, 0.60, 0.40, 24, TEAL, True, PP_ALIGN.CENTER)
    add_text(slide, "全部 collision candidates\n重新计算", 1.07, 4.35, 3.10, 0.76, 16, NAVY, True, PP_ALIGN.CENTER)
    add_box(slide, 4.91, 1.45, 4.03, 4.95, PALE_TEAL, TEAL)
    add_text(slide, "2 · Per-family gate", 5.23, 1.75, 3.2, 0.34, 20, NAVY, True)
    add_text(slide, "Eₖ(x⁺) ≤ Eₖ(x⁰) + εₖ", 5.14, 2.33, 3.56, 0.42, 22, TEAL_DARK, True, PP_ALIGN.CENTER, "Cambria Math")
    add_bullets(slide, [
        ("collision", ""), ("support", ""), ("plane", ""),
        ("semantic / depth", ""),
    ], 5.53, 3.02, 2.95, 1.62, 15)
    add_text(slide, "任一分量失败 → 不允许用 macro improvement 覆盖。", 5.35, 5.15, 3.20, 0.55, 14, RED, True, PP_ALIGN.CENTER)
    add_box(slide, 9.20, 1.45, 3.48, 4.95, WHITE, GRID)
    add_text(slide, "3 · Scoped commit", 9.50, 1.75, 2.8, 0.34, 20, NAVY, True)
    add_text(slide, "PASS", 10.07, 2.42, 1.63, 0.36, 20, GREEN, True, PP_ALIGN.CENTER)
    add_text(slide, "保留该连通组件的更新", 9.67, 2.92, 2.45, 0.55, 14, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, "LOCAL FAIL", 9.85, 3.75, 2.0, 0.36, 18, ORANGE, True, PP_ALIGN.CENTER)
    add_text(slide, "释放 witness child / factor / separator，局部重解", 9.58, 4.20, 2.56, 0.83, 13.5, DARK, False, PP_ALIGN.CENTER)
    add_text(slide, "GLOBAL FAIL", 9.85, 5.23, 2.0, 0.36, 18, RED, True, PP_ALIGN.CENTER)
    add_text(slide, "精确恢复 incumbent", 9.70, 5.67, 2.32, 0.36, 14, DARK, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 9 Results
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "Paper30 最终数据：速度比较与质量比较分开报告", "相同 v4-deepsearch S3；boundary 仍为 N/A；GT 仅事后评分", 9, total)
    runtime = [
        ["Route", "DeepSearch status", "S4 measured", "整线结论", "Speedup"],
        ["v1 / v3 · Legacy SA", "DINOv2 retrieval", "677.77s", "未做当前整线复测", "—"],
        ["v4-deepsearch · Legacy SA", "S2 w1 83.97s mean", "677.77s", "Smoke5 measured", "—"],
        ["v5 · SceneProof certified", "S2 w8 65.90s mean", "218.80s", "median 92.19s；尾延迟不稳", "S2 1.274× / S4 3.10×"],
    ]
    add_text(slide, "A · S4 3.10×；S2 8-way 平均 1.274×，但中位数退化，暂不声明整线倍率", 0.68, 1.37, 10.8, 0.32, 18, NAVY, True)
    add_table(slide, runtime, 0.62, 1.78, 12.04, 1.73, [3.2,1.8,1.8,1.8,2.2], 12.0,
              highlight_rows={3: PALE_TEAL})
    add_text(slide, "B · v5 certified 相对同源 smooth control 的质量变化", 0.68, 3.77, 6.6, 0.32, 18, NAVY, True)
    quality = [
        ["Metric", "Δ", "95% paired CI", "结论"],
        ["Physical macro (4-family)", "+0.002859", "[+0.000696, +0.005382]", "显著提升"],
        ["Collision", "+0.009823", "[+0.001475, +0.020607]", "显著提升"],
        ["Support", "−0.000117", "[−0.000878, +0.000617]", "持平"],
        ["Plane", "−0.000307", "[−0.000792, +0.000150]", "持平"],
        ["Semantic", "+0.000608", "[−0.000846, +0.002727]", "持平"],
        ["Rotation AUC@60", "+0.000150", "[−0.001054, +0.001379]", "非劣"],
        ["Translation AUC@0.5m", "+0.000083", "[−0.000528, +0.000748]", "非劣"],
    ]
    add_table(slide, quality, 0.62, 4.16, 12.04, 2.52, [3.0,1.6,3.2,1.6], 11.5,
              highlight_rows={1: PALE_TEAL, 2: PALE_TEAL})
    add_text(slide, "Smoke5 S2-only：w1 mean/median=83.974/81.358s；w8=65.904/92.187s；mean speedup=1.274×，0 failures。平均值改善但中位数变差，说明 ngrok/上游队列尾延迟明显；整线倍率必须等端到端成对计时。fix25 另发现 curtain edge-on visibility regression，当前结论标为 provisional。", 0.68, 6.74, 11.7, 0.42, 9.0, MID)
    add_footer(slide)

    # 10 Visual comparison
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "固定相机视觉对比：bedroom_01", "同视角展示 Legacy → smooth → guarded → certified；数值改善不等于所有局部关系已被证明", 10, total)
    variants = [
        ("00_legacy_sa5000.png", "Legacy SA5000", RED),
        ("01_smooth_control.png", "Smooth control", BLUE),
        ("02_guarded_raw.png", "Guarded raw", ORANGE),
        ("03_certified_final.png", "Certified final", TEAL_DARK),
    ]
    for i, (name, label, color) in enumerate(variants):
        x = 0.62 + i * 3.05
        add_picture_cover(slide, RENDERS / "bedroom_01" / name, x, 1.52, 2.82, 3.78, color, 1.7)
        add_box(slide, x, 5.39, 2.82, 0.55, color, color, radius=False)
        add_text(slide, label, x + 0.07, 5.51, 2.68, 0.28, 13, WHITE, True, PP_ALIGN.CENTER)
    add_box(slide, 0.62, 6.15, 12.03, 0.66, WHITE, GRID)
    add_runs(slide, [
        ("观察：", NAVY, True),
        ("灯具/床面布置总体更稳定；但右侧枕头仍暴露“接触但重心不稳定”的 metric blind spot。", DARK, False),
    ], 0.89, 6.31, 11.50, 0.32, 14.5)
    add_footer(slide)

    # 11 limitation
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, "已知缺口：接触成立，不代表静态与扰动稳定", "枕头 overhang 是下一版 executable SUPPORT program 的直接目标", 11, total)
    add_picture_cover(slide, support_crop, 0.62, 1.47, 5.18, 4.94, RED, 1.8)
    add_box(slide, 6.08, 1.47, 6.60, 2.24, PALE_RED, RED)
    add_text(slide, "当前 verifier 为什么漏掉", 6.40, 1.76, 4.4, 0.34, 20, RED, True)
    add_bullets(slide, [
        ("Contact gap：", "只证明表面接近；"),
        ("Footprint overlap：", "不等价于质心投影位于稳定区域；"),
        ("短时 simulation：", "passive / friction / collision proxy 可能阻止真实坠落。"),
    ], 6.40, 2.24, 5.78, 1.20, 14.5, bullet_color=RED)
    add_box(slide, 6.08, 3.98, 6.60, 2.43, PALE_TEAL, TEAL)
    add_text(slide, "下一版 SUPPORT_PATCH certificate", 6.40, 4.27, 5.2, 0.34, 20, TEAL_DARK, True)
    add_text(slide, "Πg(cCOM) ∈ Pstable   且   d(Πg(cCOM), ∂Pstable) ≥ m", 6.35, 4.83, 5.98, 0.40, 20, NAVY, True, PP_ALIGN.CENTER, "Cambria Math")
    add_bullets(slide, [
        ("Part geometry：", "bottom contact patch ↔ parent support patch；"),
        ("COM margin：", "质心投影与支撑多边形边界保持安全余量；"),
        ("Executable probe：", "小扰动 + gravity settle 后仍保持 contact topology。"),
    ], 6.40, 5.38, 5.70, 0.83, 13.5)
    add_text(slide, "这会把 SceneProof 从 object-level geometric certificate 推向真正的 functional proof。", 0.77, 6.60, 11.58, 0.36, 15, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 12 conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6]); slide.background.fill.solid(); slide.background.fill.fore_color.rgb = NAVY
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.14)); set_fill(top, TEAL); top.line.fill.background()
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.36), prs.slide_width, Inches(0.14)); set_fill(bot, TEAL); bot.line.fill.background()
    add_text(slide, "结论与下周计划", 0.88, 0.70, 6.4, 0.55, 31, WHITE, True)
    add_text(slide, "SceneProof-Core 数值门禁通过 Paper30，但 visibility certificate 尚未通过；fix25 暂为 provisional。", 0.90, 1.37, 11.2, 0.45, 17, TEAL, True)
    add_box(slide, 0.88, 2.08, 5.76, 3.92, NAVY_2, TEAL)
    add_text(slide, "本周已证明", 1.20, 2.40, 3.6, 0.34, 21, WHITE, True)
    add_bullets(slide, [
        ("S4 3.10×：", "严格实测优于 SA5000；S2 8-way 平均 1.274×，但中位数未改善。"),
        ("显著增益：", "physical macro 与 collision bootstrap CI > 0。"),
        ("Pose 非劣：", "rotation / translation 保持 recovery 主线。"),
        ("证书边界：", "211 个 changes 仅通过物理分量证书；尚不能证明对象可见性不退化。"),
    ], 1.20, 2.98, 4.90, 2.28, 15.5, color=WHITE, bullet_color=TEAL)
    add_box(slide, 6.92, 2.08, 5.54, 3.92, NAVY_2, ORANGE)
    add_text(slide, "下周只做一条主线", 7.24, 2.40, 3.8, 0.34, 21, WHITE, True)
    add_bullets(slide, [
        ("1 · Attachment axis：", "薄片厚度轴对齐墙法线，修复 curtain edge-on。"),
        ("2 · Visibility certificate：", "固定相机 object-ID area + scoped rollback。"),
        ("3 · COM certificate：", "稳定多边形 + margin。"),
        ("4 · Perturbation probe：", "独立 micro-simulation。"),
    ], 7.24, 2.98, 4.62, 2.28, 15.5, color=WHITE, bullet_color=ORANGE)
    add_text(slide, "Paper claim", 0.92, 6.35, 1.55, 0.28, 13, TEAL, True)
    add_text(slide, "关系程序定义流形、稀疏性、失败 witness 与提交证书；优化结果因而成为 proof-carrying scene。", 2.35, 6.28, 9.75, 0.42, 16, WHITE, True)
    add_text(slide, "12 / 12", 11.82, 6.95, 0.72, 0.24, 10, RGBColor(190, 203, 230), True, PP_ALIGN.RIGHT)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(OUT)


if __name__ == "__main__":
    build()
