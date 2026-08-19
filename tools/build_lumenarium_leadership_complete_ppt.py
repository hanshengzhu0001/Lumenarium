"""!!! 这个脚本已落后于 pptx，直接重跑会覆盖掉后来的改动 !!!

2026-08-18 之后的三批改动是直接落在 pptx 上的，没有回写到这里：
  1. 25 页 speaker notes 全部口语化重写（秒数改为按字数推算）；
  2. 第 4 页更正调用量口径（几万次特征比对是本地 GPU 计算，不是 API 调用）；
  3. 第 12/13 页把不可微、LM 各符号、Schur 改写成生活化说法。

所以 pptx 是唯一真源，这个脚本只保留版面与配色的生成逻辑，供重建骨架时参考。
要改内容请直接改 pptx；要重建骨架请先把上面三批改动搬回来。
"""

"""生成完整版领导汇报 PPT：在最终版 16 页基础上补齐数学、自评、拓展与致谢。

为什么另建一个脚本而不是改原脚本
--------------------------------
原脚本 build_lumenarium_final_leadership_ppt.py 产出的 16 页已经讲过一轮，
它是一个可回溯的基线。本脚本新增 6 页并重排页码；若直接改写原脚本就失去了
上一轮讲了什么的记录。两者输出到不同文件名，可并存比较。

为什么文案里的中文引号一律写成 \\u201c / \\u201d 转义
----------------------------------------------------
源码保持纯 ASCII 引号字符。曾经直接写全角引号，结果在写入环节被规范化成
ASCII 双引号，把 add_text(s, "把"单图 3D ...") 变成了提前闭合的字符串，
Python 报 invalid decimal literal。转义写法让源码不含裸全角引号，
运行时仍然产出正确的全角引号，因此这个故障不可能再次发生。

新增的 6 页只陈述能在代码里指到行号的内容
----------------------------------------
公式、权重、阈值、迭代次数均取自源码默认值，不取自记忆或推测：
  目标函数与权重       modules/_s4_layoutvlm_ops.py:3612-3629
  LM 阻尼更新规则      modules/_s4_layoutvlm_ops.py:6669-6678
  matrix-free 正规方程 modules/_s4_layoutvlm_ops.py:126-138 的 docstring
  Marquardt 缩放       modules/_s4_layoutvlm_ops.py:5065-5069
  叶子平移可消元条件    modules/_s4_scenelm_relational.py:886-888
  碰撞析取冻结的理由    modules/_s4_layoutvlm_ops.py:491-498
  接触切换处的兜底      modules/_s4_layoutvlm_ops.py:5786-5792
  最小坐标基           modules/_s4_scenelm_relational.py:359-386
  physical macro 定义  eval_physical_realizability.py:26,256-259,937-981,1143-1156
  boundary 族被排除     eval_physical_realizability.py:1009-1013,480-490
  V3 开关在 V5 仍开     run_imaginarium_I2Layout_v5_scenelm.py:11-13

被否决的表述
------------
不写 3.513x 来自二阶优化。3.513x 实测自 V5-fast/Fix61，其 solver 默认是
Adam-400（run_imaginarium_I2Layout_v4_fast.py:12-13），而二阶 solver
v5_scenelm 在 EVAL_DASHBOARD.ascii:367 仍标注 Default status OFF。
因此加速归因写成关系定向加只更新违反约束的自由度加 Schur 消元，
二阶作为求解层介绍其数学动机，不承担速度数字。
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from tools.build_lumenarium_intern_summary import (
    BLUE, DARK, GREEN, GRID, LIGHT, MID, NAVY, NAVY2, ORANGE, PALE_BLUE,
    PALE_ORANGE, PALE_RED, PALE_TEAL, PURPLE, RED, TEAL, TEAL_DARK, WHITE,
    add_arrow, add_box, add_bullets, add_footer, add_metric, add_picture_cover,
    add_runs, add_stage, add_table, add_text, add_title, set_fill,
)


ROOT = Path(__file__).resolve().parents[1]
VIS = ROOT / "visual_results"
OUT = ROOT / "Lumenarium_\u6691\u671f\u5b9e\u4e60\u603b\u7ed3_\u5b8c\u6574\u7248_2026-08-19.pptx"
NOTES = ROOT / "Lumenarium_\u6691\u671f\u5b9e\u4e60\u603b\u7ed3_\u5b8c\u6574\u7248_\u9010\u9875\u8bb2\u7a3f.md"

TOTAL = 25
MATH = "Cambria Math"
LQ = "\u201c"
RQ = "\u201d"

# 剪辑好的三档录屏放到下面任一路径（或用 LUMENARIUM_DEMO_VIDEO 指定），
# 重跑本脚本即自动嵌入到录屏页；文件不存在时该页画同尺寸的封面框，
# 因此视频到位前后版面完全一致，不会发生位移。
DEMO_VIDEO_CANDIDATES = ("demo_video.mp4", "media/lumenarium_demo.mp4")
# 三档合剪的最终录屏（Lumenarium_Final Demo）。URL 只用于显示，LINK 用于超链接。
# PowerPoint 的「联机视频」只支持 YouTube 等少数站点，B 站无法真正内嵌，
# 所以这里的做法是：本地有 mp4 就真嵌入原地播放，没有就把封面与链接文字都挂成超链接。
# 换视频只改这两行，录屏页与致谢页会同步。
DEMO_VIDEO_URL = "bilibili.com/video/BV1Ud8u6DEYE"
DEMO_VIDEO_LINK = "https://www.bilibili.com/video/BV1Ud8u6DEYE/"


def q(text: str) -> str:
    """Wrap text in full-width quotes without writing them literally."""
    return LQ + text + RQ


# 幻灯片的构建顺序与放映顺序刻意分开。
# 构建顺序按代码里的历史分块（封面、痛点、录屏、演示、优先级、继承……），
# 放映顺序则按叙事需要重排：先把问题与继承的系统讲透（含 S0–S4 全链路），
# 再放录屏、再现场跑，这样演示不显得无来由。
# 下面这个元组读作「放映第 k 页 = 构建顺序里的第 SLIDE_ORDER[k-1] 页」。
# 改顺序只改这一处：页码、备注编号、讲稿顺序都从它派生，不会各自走偏。
SLIDE_ORDER = (1, 2, 5, 6, 7, 9, 8, 3, 4, 10, 11, 12, 13, 14, 15,
               16, 17, 18, 19, 20, 21, 22, 23, 24, 25)

_PAGE = {"n": 1}  # slide 1 is the cover and draws no page number


def page() -> int:
    """返回当前正在构建的这一页在放映顺序里的位置。

    页码必须是听众看到的顺序，而不是代码里的书写顺序，所以这里做一次翻译。
    """
    _PAGE["n"] += 1
    return SLIDE_ORDER.index(_PAGE["n"]) + 1


def reorder(prs) -> None:
    """把幻灯片重排成放映顺序。

    在这里重排、而不是搬动几百行建页代码，是为了让「顺序」成为一个可以一行改动、
    一眼看懂的决定；建页代码保持原样也便于和历史版本对照。
    """
    id_list = prs.slides._sldIdLst
    built = list(id_list)
    if len(built) != len(SLIDE_ORDER):
        raise SystemExit(f"建了 {len(built)} 页，但 SLIDE_ORDER 写了 {len(SLIDE_ORDER)} 页")
    for element in built:
        id_list.remove(element)
    for index in SLIDE_ORDER:
        id_list.append(built[index - 1])


def base_slide(prs: Presentation, title: str, subtitle: str, page_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_title(slide, title, subtitle, page_no, TOTAL)
    add_footer(slide, "Lumenarium \u00b7 \u6691\u671f\u5b9e\u4e60\u603b\u7ed3 \u00b7 Hansen Zhu \u00b7 2026-08")
    return slide


def add_four_up(slide, files, labels, y=1.55, h=3.92):
    xs = [0.55, 3.74, 6.93, 10.12]
    for x, file, label in zip(xs, files, labels):
        add_picture_cover(slide, file, x, y, 2.72, h)
        add_box(slide, x + 0.10, y + h - 0.54, 2.52, 0.42, NAVY, NAVY)
        add_text(slide, label, x + 0.17, y + h - 0.44, 2.38, 0.22,
                 11.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def link_shape(shape, url: str = DEMO_VIDEO_LINK):
    """让放映时点击这个形状即跳转到视频页。

    只挂 click_action、不动文字，因为超链接文字会被 PowerPoint 改成主题色并加下划线，
    而封面上那块牌子的白字是版面的一部分。失败只是少一次跳转，不能让构建挂掉。
    """
    try:
        shape.click_action.hyperlink.address = url
    except Exception:
        pass
    return shape


def link_text(box, url: str = DEMO_VIDEO_LINK):
    """把一行文字变成可点击链接；下划线是刻意保留的，它提示这行可以点。"""
    try:
        for paragraph in box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.hyperlink.address = url
    except Exception:
        pass
    return box


def find_demo_video() -> Path | None:
    override = os.environ.get("LUMENARIUM_DEMO_VIDEO")
    candidates = [Path(override)] if override else []
    candidates += [ROOT / name for name in DEMO_VIDEO_CANDIDATES]
    return next((path for path in candidates if path.is_file()), None)


def add_demo_video(slide, x, y, w, h, poster: Path) -> bool:
    """嵌入剪辑好的录屏；文件还不存在时画同尺寸的封面框。

    这一页必须在视频剪好之前就能生成，所以缺文件是正常状态而不是错误。
    退化路径刻意用完全相同的位置与尺寸，视频到位后重跑脚本，其余元素不会移动。
    add_movie 失败也走同一条退化路径：宁可少一段视频，不能在汇报前一晚构建失败。
    """
    video = find_demo_video()
    if video is not None:
        try:
            slide.shapes.add_movie(
                str(video), Inches(x), Inches(y), Inches(w), Inches(h),
                poster_frame_image=str(poster), mime_type="video/mp4",
            )
            return True
        except Exception:
            pass
    add_picture_cover(slide, poster, x, y, w, h)
    link_shape(add_box(slide, x + w / 2 - 1.55, y + h / 2 - .50, 3.10, 1.00, NAVY, TEAL, width=1.6))
    link_shape(add_text(slide, "\u25b6  播放三档实机录屏",
                        x + w / 2 - 1.45, y + h / 2 - .22, 2.90, .40, 15, WHITE, True,
                        PP_ALIGN.CENTER))
    return False


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Lumenarium 暑期实习总结（完整版）"
    prs.core_properties.subject = "从单图重建到更快、可证明、可部署的 3D 场景系统"
    prs.core_properties.author = "Hansen Zhu"

    living = [VIS / "livingroom_10_worst_input.png", VIS / "livingroom_10_worst_v1_final.png",
              VIS / "livingroom_10_worst_v3_final.png", VIS / "livingroom_v5_final.png"]
    office = [VIS / "official_02_worst_input.png", VIS / "official_02_worst_v1_final.png",
              VIS / "official_02_worst_v3_final.png", VIS / "official_v5_final.png"]

    # ---------------------------------------------------------------- 1 cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.13))
    set_fill(band, TEAL); band.line.fill.background()
    add_picture_cover(s, living[3], 8.25, 0, 5.08, 7.5, border=NAVY, border_width=0)
    add_text(s, "SUMMER INTERNSHIP REVIEW \u00b7 2026", .78, .67, 6.5, .3, 13, TEAL, True)
    add_text(s, "Lumenarium", .78, 1.43, 7.0, .72, 42, WHITE, True)
    add_text(s, "从一张图，到更快、可验证、可交付的 3D 场景", .80, 2.28, 7.20, 1.10, 27, WHITE, True)
    add_box(s, .80, 4.00, 6.36, 1.18, NAVY2, TEAL)
    add_runs(s, [("3.513\u00d7 S4 提速  ", TEAL, True), ("+7.52 pp physical macro", WHITE, False)],
             1.08, 4.31, 5.82, .5, 17, MSO_ANCHOR.MIDDLE)
    add_text(s, "Hansen Zhu \u00b7 Mentor: Calvin Gu \u00b7 腾讯暑期实习", .81, 6.40, 6.3, .28, 13, WHITE, True)

    # ------------------------------------------------------- 2 一个真实的痛点
    # 页码一律由 page() 发放；下面各段注释里的序号是历史顺序，插页后不再等于页码。
    s = base_slide(prs, "先讲一个场景：技术美术拿到一张参考图",
                   "他要的不是" + q("像") + "，而是一个能直接开工的三维场景", page())
    add_box(s, .62, 1.40, 5.86, 2.30, PALE_RED, RED, width=1.5)
    add_text(s, "他会怎么用", .95, 1.70, 2.4, .34, 19, RED, True)
    add_bullets(s, ["导进 Blender / UE，摆机位打光",
                    "替换个别资产，微调布局",
                    "交给下游做动画、碰撞、物理"],
                .96, 2.26, 5.00, 1.30, 14, bullet_color=RED)
    add_box(s, .62, 3.92, 5.86, 2.42, WHITE, GRID)
    add_text(s, "所以一个" + q("看起来对") + "的场景，可能完全不能用",
             .95, 4.22, 5.45, .34, 18, NAVY, True)
    add_bullets(s, [("桌子穿进墙里：", "导进引擎第一帧就炸"),
                    ("杯子浮在桌面上方两厘米：", "渲染图看不出，动画一动就露"),
                    ("窗帘没被检索出来：", "整个镜头要重做")],
                .96, 4.80, 5.00, 1.42, 13, bullet_color=RED)
    add_box(s, 6.80, 1.40, 5.90, 4.94, NAVY, NAVY)
    add_text(s, "这两个月我在回答的问题", 7.14, 1.78, 4.6, .36, 20, TEAL, True)
    add_text(s, "怎么让系统\n只提交它能证明的修改？", 7.14, 2.52, 5.20, 1.40, 30, WHITE, True)
    add_box(s, 7.14, 4.20, 5.20, 1.86, NAVY2, TEAL)
    add_text(s, "不是让它搜索更多可能，", 7.42, 4.52, 4.70, .34, 16, WHITE, False)
    add_text(s, "而是让每一次改动都带着证据落地。", 7.42, 4.94, 4.70, .34, 16, TEAL, True)
    add_text(s, "这句话决定了后面所有技术选择。", 7.42, 5.48, 4.70, .30, 13, RGBColor(199, 219, 238))

    # ---------------------------------------------- 3 三档实机录屏（本次新增）
    s = base_slide(prs, "先看它做完之后是什么样：三档实机录屏",
                   "V5-fast / V5-medium / V5-best，同一张输入三种交付口径。剪辑自真实运行，倍速播放", page())
    add_demo_video(s, .62, 1.40, 7.60, 4.28, office[3])
    profiles = [(TEAL, "V5-fast", "冻结 Fix61，不做任何展示性修改",
                 "论文定量口径：只有它可以进主表"),
                (BLUE, "V5-medium", "在此之上加保守的 visual-safe 清理",
                 "演示与技术美术交付"),
                (PURPLE, "V5-best", "再加全对象真实支撑审计与事务化 first-contact 掉落",
                 "最高物理完整性；证不了的报 unresolved")]
    for index, (color, name, what, who) in enumerate(profiles):
        y = 1.40 + index * 1.46
        add_box(s, 8.42, y, 4.28, 1.36, WHITE, color, width=1.6)
        add_text(s, name, 8.68, y + .10, 2.6, .28, 15, color, True)
        add_text(s, what, 8.68, y + .44, 3.78, .50, 11.6, DARK)
        add_text(s, who, 8.68, y + .98, 3.78, .28, 11.4, MID, True)
    add_box(s, .62, 5.84, 12.08, 1.02, NAVY, NAVY)
    add_text(s, "三档不是三个模型", .92, 6.00, 2.6, .26, 13, TEAL, True)
    add_text(s, "它们是同一条流水线上三种不同的提交策略：越往右提交越保守、要求的证据越多，"
                "而 S0\u2013S3 的重建结果共享同一份冻结缓存。所以换档不重跑重建，只重跑最后一步的策略。",
             3.10, 5.96, 9.32, .56, 12.2, WHITE)
    link_text(add_text(s, "完整录屏：" + DEMO_VIDEO_URL,
                       3.10, 6.54, 9.32, .24, 11.4, TEAL, True))

    # ------------------------------------------------- 4 实机演示（本次新增）
    s = base_slide(prs, "那就现在跑一次：同一条链路，现场从头走",
                   "上传一张办公室参考图，固定随机种子，从头把五个阶段全量跑一遍", page())
    # 这一页只放听众需要知道的事实。操作顺序、怎么判读进度、网络不通怎么办，
    # 都属于我自己的提示，写在备注里；把提示投到屏幕上只会分散注意力。
    add_picture_cover(s, office[0], .62, 1.40, 4.26, 4.26)
    add_box(s, .72, 5.12, 4.06, .44, NAVY, NAVY)
    add_text(s, "今天的输入 \u00b7 1024 \u00d7 1024 参考图",
             .78, 5.20, 3.94, .28, 11.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_box(s, 5.10, 1.40, 7.60, 1.76, NAVY, NAVY)
    add_text(s, "现在打开这个地址", 5.40, 1.58, 4.0, .26, 13.5, TEAL, True)
    add_text(s, "embedding.lightart.qq.com", 5.40, 1.90, 7.0, .60, 29, WHITE, True)
    add_text(s, "跑在两块 A10 上的在线服务：给它一张室内参考图，拿到能直接在 Blender 或 UE 里"
                "打开的场景。",
             5.40, 2.56, 7.02, .52, 12.4, RGBColor(202, 220, 239))
    add_box(s, 5.10, 3.30, 7.60, .96, PALE_TEAL, TEAL, width=1.6)
    add_text(s, "今天这次怎么跑", 5.40, 3.44, 5.0, .26, 13.5, TEAL_DARK, True)
    add_text(s, "同一张图我们会缓存它的重建结果，第二次提交几十秒就能出图。今天不用缓存："
                "固定随机种子，从头把五个阶段全量跑一遍。",
             5.40, 3.76, 7.02, .44, 12.4, DARK)
    steps = [(TEAL, "1", "上传参考图", "就是左边这张。"),
             (BLUE, "2", "选 V5-demo", "跑最完整的那条流水线，随机种子固定住。"),
             (PURPLE, "3", "Generate scene", "五个阶段依次点亮，大约十几分钟。")]
    for index, (color, number, head, body) in enumerate(steps):
        x = 5.10 + index * 2.58
        add_box(s, x, 4.40, 2.44, 1.48, WHITE, color, width=1.6)
        add_text(s, number, x + .20, 4.50, .42, .38, 21, color, True)
        add_text(s, head, x + .66, 4.56, 1.72, .28, 13.5, NAVY, True)
        add_text(s, body, x + .22, 4.92, 2.02, .90, 11.5, DARK)
    add_box(s, .62, 6.06, 12.08, .84, NAVY, NAVY)
    add_text(s, "输出稍后在浏览器里看", .92, 6.22, 3.0, .26, 13, TEAL, True)
    add_text(s, "这一页只放输入。它现在在后台真跑，我们接着往下讲；跑完之后我直接切回浏览器，"
                "看这一次的场景和它实际用的随机种子。",
             3.60, 6.20, 8.82, .56, 12.4, WHITE)

    # ----------------------------------------------------------- 5 priorities
    s = base_slide(prs, "先定义" + q("好场景") + "：光子真正需要什么？",
                   "不是所有指标同等重要；两个月的所有技术选择都围绕这个排序", page())
    add_box(s, .62, 1.40, 7.22, 4.95, WHITE, TEAL, width=1.7)
    add_text(s, "交付优先级", .95, 1.72, 2.5, .36, 21, TEAL_DARK, True)
    priorities = [("01  资产完整性", "主体与关键道具不能缺席"),
                  ("02  速度", "能进入真实制作链路，而不是离线实验"),
                  ("03  不穿模", "最影响可信度的三维错误"),
                  ("04  物体间关系", "桌上、墙上、床上等结构要成立")]
    for i, (head, body) in enumerate(priorities):
        y = 2.31 + i * .82
        add_text(s, head, 1.00, y, 2.62, .30, 16, NAVY, True)
        add_text(s, body, 3.52, y, 3.74, .31, 14, DARK)
    add_box(s, 8.15, 1.40, 4.56, 4.95, PALE_ORANGE, ORANGE, width=1.5)
    add_text(s, "可以接受的 trade-off", 8.48, 1.72, 3.60, .36, 20, ORANGE, True)
    add_bullets(s, [("单体资产 pose：", "rotation / translation 可略有牺牲。"),
                    ("少量非关键悬空：", "优先显式 unresolved，不伪装成功。")],
                8.48, 2.42, 3.48, 1.54, 14, bullet_color=ORANGE)
    add_box(s, 8.48, 4.44, 3.50, 1.20, NAVY, NAVY)
    add_text(s, "核心判断", 8.73, 4.66, 1.2, .26, 14, TEAL, True)
    add_text(s, "场景整体可用性 > 单个资产像素级复刻", 8.73, 5.00, 2.98, .60, 14, WHITE, True)

    # ------------------------------------------- 4 我们继承的 Imaginarium
    s = base_slide(prs, "我们继承的系统：Imaginarium 把路走通了",
                   "四个阶段都能跑通，但每一步的代价最后都记在同一个地方", page())
    add_box(s, .62, 1.34, 4.06, 2.42, WHITE, TEAL, width=1.6)
    add_text(s, "资产是怎么被选中并放进场景的", .90, 1.56, 3.6, .30, 15, TEAL_DARK, True)
    add_bullets(s, [("粗排：", "DINOv2 patch 特征，把物体切图与资产 500+ 视角标准照比相似度"),
                    ("精排：", "VLM 逐个比对纹理与尺寸，挑语义最优"),
                    ("落位：", "几何与视觉双通道打分选定资产，再用 RANSAC 把 3D 框对齐到 2D 观测")],
                .91, 1.94, 3.50, 1.74, 11.4, bullet_color=TEAL)
    add_box(s, 4.84, 1.34, 3.94, 2.42, PALE_RED, RED, width=1.6)
    add_text(s, "为什么 API 调用那么多", 5.12, 1.56, 3.4, .30, 15, RED, True)
    add_text(s, "检索复杂度 = 物体数 \u00d7 视角数 \u00d7 算子数", 5.13, 1.94, 3.50, .28, 12, NAVY, True)
    add_text(s, "约 50 物体 \u00d7 500+ 视角 \u2248 2.5 万次比对；其中每个候选还要 2\u20135 次 VLM 调用做纹理精排。"
                "真正的瓶颈不是次数，而是限流：有效并发只有 1\u20132，全部串行等待，"
                "单场景 60\u2013180 s 就耗在这里。",
             5.13, 2.28, 3.50, 1.40, 11.4, DARK)
    add_box(s, 8.94, 1.34, 3.76, 2.42, PALE_ORANGE, ORANGE, width=1.6)
    add_text(s, "支撑关系由 GPT 以文字给出", 9.22, 1.56, 3.2, .30, 15, ORANGE, True)
    add_text(s, "检测是视觉的（GroundingDINO / SAM3 出框和 mask），但" + q("谁在谁上面") +
                "这个关系，由 GPT 以文本 JSON 给出 supported 字段。它会系统性判错："
                "本该在衣柜顶上的手提箱被标成 floor_0，于是物理模拟时它从柜顶掉下并穿透地面。"
                "一个用语言回答的几何问题。",
             9.23, 1.94, 3.30, 1.74, 11.4, DARK)
    add_box(s, .62, 3.92, 6.02, 1.34, WHITE, PURPLE, width=1.6)
    add_text(s, "为什么最后要退火 5,000 步", .90, 4.12, 3.6, .30, 15, PURPLE, True)
    add_text(s, "每个物体只保留 x、y、\u03b8 三个自由度，随机扰动加 Metropolis 接受、指数降温。"
                "它没有梯度方向，只能用迭代次数换收敛概率。代价是每轮 N\u00b2 次两两重叠检测，"
                "20 物体乘上万轮就是数百万次判定。",
             .92, 4.48, 5.44, .74, 11.6, DARK)
    add_box(s, 6.82, 3.92, 5.88, 1.34, NAVY, NAVY)
    add_text(s, "一句话", 7.10, 4.12, 1.6, .30, 15, TEAL, True)
    add_text(s, "前面每一步不确定，最后一步就得用更多搜索去兜。"
                "所以我没有从渲染质量入手，而是从这条因果链的上游入手。",
             7.12, 4.48, 5.32, .74, 13.5, WHITE, True)
    reuse = [["我们的取舍", "内容", "理由"],
             ["复用", "视觉引导重建、资产库、场景布局表示、深度与 OBB 先验",
              "这套骨架是有效的，没有理由重写"],
             ["重做", "检索、支撑关系判定、最后一步的优化与提交",
              "这三处正是误差被放大的地方"],
             ["刻意不做", "Flux 微调（论文借它把概念图分布对齐资产库，Top-1 检索 48.6% \u2192 68.7%）",
              "它是对某一批资产过拟合；我们要能换场景换资产库，不能把输入分布锁死"]]
    add_table(s, reuse, .62, 5.40, 12.08, 1.02, [1.40, 6.86, 3.82], 10.6, {3: PALE_TEAL})

    # ------------------------------------------------------------ 5 V1 痛点
    s = base_slide(prs, "原始系统不是" + q("不能生成") + "，而是难以稳定交付",
                   "看起来像一个房间，不等于三维关系真的成立", page())
    add_picture_cover(s, living[1], .62, 1.38, 6.05, 4.90)
    add_box(s, 6.94, 1.38, 5.76, 4.90, WHITE, GRID)
    add_text(s, "四类阻碍", 7.27, 1.72, 2.0, .34, 20, NAVY, True)
    add_bullets(s, [("不完整：", "窗帘、支撑物或关键对象在链路中丢失"),
                    ("太慢：", "S4 的 SA-5000 单场景约 678 秒"),
                    ("会穿模：", "局部视觉合理，但真实网格已重叠"),
                    ("关系不可信：", "错误 parent、悬空和错误上墙会级联")],
                7.28, 2.25, 4.70, 2.60, 14.0, bullet_color=RED)
    add_box(s, 7.24, 5.20, 4.86, .78, PALE_TEAL, TEAL)
    add_text(s, "因此目标不是" + q("多搜索") + "，而是" + q("只提交能证明的修改"),
             7.46, 5.37, 4.44, .46, 12.4, TEAL_DARK, True, PP_ALIGN.CENTER)

    # -------------------------------------------------------------- 5 journey
    s = base_slide(prs, "两个月的主线：结构 \u2192 速度 \u2192 证明 \u2192 服务",
                   "每一个失败样例都变成了下一层系统能力", page())
    milestones = [("V1", "Baseline", "Imaginarium 重建", BLUE), ("V3", "Support tree", "支撑 / stack", ORANGE),
                  ("V4", "DeepSearch", "更快检索", PURPLE), ("V5", "LM + Proof", "更快、更物理", TEAL),
                  ("API", "Service", "双 A10 交付", GREEN)]
    xs = [.58, 3.12, 5.66, 8.20, 10.74]
    for i, (code, title, detail, color) in enumerate(milestones):
        add_stage(s, xs[i], 2.08, 2.03, 2.35, code, title, detail, color)
        if i < 4:
            add_arrow(s, xs[i] + 2.06, 3.22, xs[i + 1] - .06, 3.22, MID, 1.8)
    add_box(s, .78, 5.08, 11.76, 1.10, NAVY, NAVY)
    add_text(s, "一句话：从" + q("全局盲搜一个看起来不错的答案") + "，转向"
             + q("关系定向修改 + 几何证书提交") + "。",
             1.13, 5.39, 11.05, .38, 18, WHITE, True, PP_ALIGN.CENTER)

    # ------------------------------------------------------------- 6 pipeline
    s = base_slide(prs, "当前 Lumenarium：一张图走完整 S0\u2013S4",
                   "继承 Imaginarium 的视觉重建基础，重做检索、优化、证明和部署", page())
    stages = [("S0", "Geometry", "深度 / 相机", BLUE), ("S1", "Understand", "SAM3 / 场景图", TEAL),
              ("S2", "Retrieve", "DeepSearch / DINOv2", ORANGE), ("S3", "Recover", "pose / support tree", PURPLE),
              ("S4", "Optimize", "SceneLM / SceneProof", GREEN)]
    # 每一阶段都写清「吃什么 → 吐什么」：这条链路是后面所有痛点与改进的挂载点，
    # 只画五个方框而不写接口，听众就无法判断某个改进到底动了哪一段。
    flow = [(BLUE, "吃：一张 1024\u00d71024 概念图",
             "吐：单目度量深度图、相机内参、点云"),
            (TEAL, "吃：概念图 + 深度",
             "吐：实例 mask 与类别、有向包围盒 OBB、场景图（谁被谁支撑）"),
            (ORANGE, "吃：每个实例的切图 + OBB",
             "吐：资产库里的 Top-K 候选模型"),
            (PURPLE, "吃：候选资产 + OBB + mask",
             "吐：每个物体的 6D 位姿与缩放、支撑树"),
            (GREEN, "吃：全部 6D 位姿 + 支撑树",
             "吐：无碰撞布局、渲染图、几何证书")]
    xs = [.56, 3.09, 5.62, 8.15, 10.68]
    for i, item in enumerate(stages):
        add_stage(s, xs[i], 1.42, 2.08, 2.10, *item)
        if i < 4:
            add_arrow(s, xs[i] + 2.10, 2.36, xs[i + 1] - .05, 2.36, MID, 1.8)
        color, eats, emits = flow[i]
        add_box(s, xs[i], 3.62, 2.08, 1.72, WHITE, color, width=1.2)
        add_text(s, eats, xs[i] + .16, 3.76, 1.76, .60, 9.6, MID)
        add_text(s, emits, xs[i] + .16, 4.42, 1.76, .84, 9.6, DARK, True)
    add_box(s, .56, 5.50, 12.20, 1.34, NAVY, NAVY)
    add_runs(s, [("SceneLM  ", TEAL, True), ("关系范围内定向求解", WHITE, False),
                 ("     SceneProof  ", TEAL, True), ("证书 \u2192 commit / restore / unresolved", WHITE, False)],
             .84, 5.76, 11.64, .44, 15.5, MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    add_text(s, "输出：editable placement \u00b7 render \u00b7 evaluation certificate \u00b7 result bundle",
             .91, 6.32, 11.50, .25, 11, RGBColor(199, 219, 238), False, PP_ALIGN.CENTER)

    # --------------------------------------------------- 7 V3 -> V5 support 线
    s = base_slide(prs, "V3 的 support-aware：用几何判据替掉文字判断",
                   "这条线 S3 侧被 V5 完整继承，S4 侧则被关系程序替换", page())
    add_box(s, .62, 1.38, 4.02, 4.96, PALE_TEAL, TEAL, width=1.6)
    add_text(s, "support-aware 三步做法", .92, 1.66, 3.5, .34, 18, TEAL_DARK, True)
    # 先分清「检测 vs 关系」：检测是视觉模型在做，关系曾是 GPT 的文字，这里换成纯几何。
    # 这句因果不写出来，听众会以为"几何"是某个检测模型判的（它真不是）。
    add_text(s, "检测是视觉的（出框、mask）；关系曾是 GPT 的文字，这里换成纯几何。",
             .93, 2.02, 3.40, .44, 10.5, DARK)
    # 这一页要回答的是「怎么做」，不是「保留了什么」：三步写在幻灯片上，
    # 听众才跟得上后面「判对之后怎么摆」这个转折。
    add_bullets(s, [("① 检测堆叠对：", "水平投影重叠 + 竖直间隙落容差内，即上下关系；"
                                      "阈值全由物体尺寸算出，无硬编码"),
                    ("② 抬升并冻结：", "先算上层高度；退火期间把它从扰动集合中排除，"
                                      "退火后按下层顶面精确落位"),
                    ("③ 沿支撑链递归：", "父物体一动，后代跟随")],
                .93, 2.50, 3.34, 2.60, 11.4, bullet_color=TEAL)
    add_text(s, "仍保留：地面父验证 \u00b7 OBB 邻近父验证（Algorithm 2）\u00b7 tree_sons",
             .93, 5.62, 3.40, .28, 10, MID, True)
    add_box(s, 4.80, 1.38, 3.82, 4.96, PALE_ORANGE, ORANGE, width=1.6)
    add_text(s, "被替换掉的", 5.10, 1.68, 3.0, .34, 18, ORANGE, True)
    add_text(s, "V3 设想的 S4 堆叠感知模拟退火", 5.11, 2.22, 3.30, .56, 14, NAVY, True)
    add_text(s, "给堆叠物体的上层加 z 向扰动、加堆叠分离惩罚。这条路线在 V5 没有继续，因为它仍然是"
             + q("扰动后看看好不好") + "，无法说明为什么该动。",
             5.11, 2.90, 3.24, 1.66, 12.4, DARK)
    add_box(s, 5.11, 4.66, 3.24, 1.42, WHITE, ORANGE)
    add_text(s, "V5 的做法", 5.34, 4.86, 2.0, .26, 12, ORANGE, True)
    add_text(s, "把" + q("甲在乙上面") + "直接编译成一条可求解的约束，而不是让退火去碰运气。",
             5.34, 5.20, 2.80, .82, 12, DARK, True)
    add_box(s, 8.78, 1.38, 3.92, 4.96, WHITE, PURPLE, width=1.6)
    add_text(s, "V5 新增的一层", 9.08, 1.68, 3.0, .34, 18, PURPLE, True)
    add_bullets(s, [("真实网格质心：", "不再用包围盒近似判稳定"),
                    ("见证支撑区域：", "只在真的量到接触时才敢改"),
                    ("分量级证书：", "每个物理分量单独门控"),
                    ("事务化沉降：", "失败就精确回滚")],
                9.09, 2.22, 3.24, 2.90, 12.2, bullet_color=PURPLE)
    add_box(s, 9.09, 5.30, 3.24, .80, NAVY, NAVY)
    add_text(s, "V3 回答" + q("谁在谁上面") + "；V5 回答" + q("这样放到底稳不稳") + "。",
             9.30, 5.50, 2.86, .46, 11.6, WHITE, True, PP_ALIGN.CENTER)

    # ------------------------------------------------- 8 SceneLM 目标函数
    s = base_slide(prs, "SceneLM 到底在最小化什么？",
                   "把" + q("场景合理") + "翻译成一个带权残差平方和，每一项都能单独量出来", page())
    add_box(s, .62, 1.36, 12.08, 1.66, WHITE, TEAL, width=1.7)
    add_text(s, "总目标", .92, 1.56, 1.6, .30, 15, TEAL_DARK, True)
    add_text(s, "E(x) = 1.0\u00b7collision + 2.0\u00b7contact + 2.0\u00b7plane + 0.25\u00b7orientation",
             1.02, 1.92, 11.4, .38, 20, NAVY, True, PP_ALIGN.CENTER, font=MATH)
    add_text(s, "+ 1.0\u00b7containment + 0.5\u00b7semantic + 1.0\u00b7boundary + 1.0\u00b7depth + 0.01\u00b7warm-start",
             1.02, 2.38, 11.4, .38, 17.5, NAVY, True, PP_ALIGN.CENTER, font=MATH)
    data = [["项", "它在问什么", "权重", "读法"],
            ["contact / plane", "该贴的面贴上了吗（桌面、墙面）", "2.0", "最高：几何硬约束"],
            ["collision", "真实网格有没有重叠", "1.0", "基准权重"],
            ["containment", "抽屉里的东西还在抽屉里吗", "1.0", "基准权重"],
            ["depth", "投影回原图还对得上吗", "1.0", "锚住原图证据"],
            ["semantic", "朝向、对齐这类语义偏好", "0.5", "只有硬约束的 1/4"],
            ["orientation", "平面朝向偏差", "0.25", "弱引导"],
            ["warm-start", "别跑离上游初值太远", "0.01", "1/200：信任但不锁死"]]
    add_table(s, data, .62, 3.16, 7.60, 3.10, [1.62, 3.16, .78, 2.04], 11.4,
              {1: PALE_TEAL, 7: PALE_BLUE})
    add_box(s, 8.44, 3.16, 4.26, 3.10, PALE_BLUE, BLUE, width=1.5)
    add_text(s, "权衡是显式写下来的", 8.74, 3.44, 3.4, .32, 17, BLUE, True)
    add_text(s, "几何硬约束的权重是语义偏好的 4 倍：一把椅子朝向不完美可以接受，"
                "但椅子腿浮在地板上方不行。",
             8.75, 3.90, 3.68, 1.10, 13, DARK)
    add_text(s, "warm-start 只有 0.01，是全场最小的一项：我们信任上游给的初值，"
                "但绝不让它把答案锁死。",
             8.75, 5.06, 3.68, 1.02, 13, DARK)

    # ------------------------------------------------- 9 为什么用二阶
    s = base_slide(prs, "为什么敢用二阶方法：先把不可微的地方修成可微",
                   "一阶只知道往哪走，二阶知道走多远。但前提是这个曲面不能在脚下跳变", page())
    add_box(s, .62, 1.36, 3.94, 2.34, PALE_RED, RED, width=1.5)
    add_text(s, "一阶的困境", .92, 1.62, 2.6, .32, 17, RED, True)
    add_text(s, "SA-5000 / Adam-400", .93, 2.04, 3.3, .30, 14, NAVY, True)
    add_text(s, "梯度只给方向，步长靠调。碰撞约束本质是从若干个分离方向里选一个，"
                "是析取而非光滑函数；最小轴一跳，线性化就震荡。",
             .93, 2.44, 3.34, 1.16, 12.2, DARK)
    add_box(s, 4.72, 1.36, 3.94, 2.34, PALE_TEAL, TEAL, width=1.5)
    add_text(s, "我们做的第一步", 5.02, 1.62, 2.8, .32, 17, TEAL_DARK, True)
    add_text(s, "冻结最小平移分离方向", 5.03, 2.04, 3.3, .30, 14, NAVY, True)
    add_text(s, "一旦判定某对物体已穿透，就固定它最省力的那个分离方向，把析取变成一个稳定的"
                "标量半空间残差。这才让 Gauss\u2013Newton 的线性化不再跳。",
             5.03, 2.44, 3.34, 1.20, 12.2, DARK)
    add_box(s, 8.82, 1.36, 3.88, 2.34, WHITE, PURPLE, width=1.5)
    add_text(s, "然后才是二阶", 9.12, 1.62, 2.6, .32, 17, PURPLE, True)
    add_text(s, "Levenberg\u2013Marquardt", 9.13, 2.04, 3.3, .30, 14, NAVY, True)
    add_text(s, "用残差的曲率信息一步解到局部二次模型的谷底，而不是靠学习率试探。",
             9.13, 2.44, 3.28, .90, 12.2, DARK)
    add_box(s, .62, 3.92, 6.02, 2.42, WHITE, GRID)
    add_text(s, "每步解的方程", .92, 4.16, 3.0, .30, 16, NAVY, True)
    add_text(s, "( J\u1d40J + \u03bb\u00b7diag(J\u1d40J) ) \u03b4 = \u2212J\u1d40r",
             1.00, 4.58, 5.30, .44, 21, TEAL_DARK, True, PP_ALIGN.CENTER, font=MATH)
    add_text(s, "关键在于从不显式构造 J、J\u1d40J 或 Hessian：J\u1d40J\u00b7v 通过 JVP/VJP 两次自动微分求得，"
                "外层用 PCG 迭代解。内存与一阶同阶，却拿到二阶的收敛性。",
             .96, 5.16, 5.36, 1.02, 12.4, DARK)
    add_box(s, 6.82, 3.92, 5.88, 2.42, PALE_ORANGE, ORANGE, width=1.5)
    add_text(s, "\u03bb 自己会调：这就是有效的来源", 7.12, 4.16, 4.8, .30, 16, ORANGE, True)
    add_text(s, "\u03c1 = 实际下降 / 二次模型预测下降", 7.14, 4.58, 5.2, .32, 15, NAVY, True, font=MATH)
    add_bullets(s, [("\u03c1 > 0.75：", "模型很准 \u2192 \u03bb 减半，更像纯二阶，步子放大"),
                    ("\u03c1 < 0.25：", "模型不准 \u2192 \u03bb 加倍，退回梯度下降"),
                    ("步被拒绝：", "\u03bb \u00d74，缩小信赖域重试")],
                7.14, 5.00, 5.30, 1.24, 12, bullet_color=ORANGE)

    # ------------------------------------------- 10 稀疏结构与诚实边界
    s = base_slide(prs, "为什么快：只解该解的那部分变量",
                   "关系不只提供约束，它还决定了哪些自由度根本不需要出现在方程里", page())
    add_box(s, .62, 1.38, 6.02, 2.62, WHITE, TEAL, width=1.6)
    add_text(s, "最小坐标：关系决定自由度个数", .92, 1.64, 4.6, .32, 17, TEAL_DARK, True)
    coord = [["关系", "该物体的平移自由度"],
             ["放在支撑面上", "2（沿支撑面世界 XY）"],
             ["贴在墙 / 天花板上", "2（沿该平面切向）"],
             ["锚定在平面法向", "1（只沿法向）"],
             ["自由物体", "3"]]
    add_table(s, coord, .92, 2.10, 5.42, 1.66, [2.60, 2.82], 11.4, {1: PALE_TEAL})
    add_text(s, "尺度全程冻结；旋转永不被消元，避免父物体的 yaw 误差传染给子物体。",
             .94, 3.62, 5.40, .30, 11.6, MID, True)
    add_box(s, 6.82, 1.38, 5.88, 2.62, PALE_BLUE, BLUE, width=1.6)
    add_text(s, "Schur 消元：先解根，再回代叶子", 7.12, 1.64, 4.6, .32, 17, BLUE, True)
    add_text(s, "[ H\u1d63\u1d63  H\u1d63\u2097 ] [\u03b4\u1d63]        [ g\u1d63 ]\n"
                "[ H\u2097\u1d63  H\u2097\u2097 ] [\u03b4\u2097]  =  \u2212 [ g\u2097 ]",
             7.16, 2.10, 5.24, .84, 15, NAVY, True, PP_ALIGN.CENTER, font=MATH)
    add_text(s, "只有有父节点且自身没有子节点的叶子平移才允许被消去，且该因子的分隔集必须恰为其父。"
                "对当前线性化子问题这是精确消元，不是近似删变量。",
             7.14, 3.06, 5.28, .92, 12.2, DARK)
    add_box(s, .62, 4.24, 6.02, 2.10, PALE_ORANGE, ORANGE, width=1.5)
    add_text(s, "二阶模型会失效的地方，我们也写下来了", .92, 4.50, 5.2, .30, 16, ORANGE, True)
    add_text(s, "在接触状态切换的瞬间，光滑 Gauss\u2013Newton 模型不可靠。当 Schur 步在正反两个方向上都让"
                "真实目标上升时，退化为在已审计的责任子空间内做无导数正交搜索。它慢，"
                "但只在例外情形触发。",
             .94, 4.92, 5.44, 1.32, 12.2, DARK)
    add_box(s, 6.82, 4.24, 5.88, 2.10, NAVY, NAVY)
    add_text(s, "一句话总结这三页", 7.12, 4.52, 4.0, .30, 16, TEAL, True)
    add_text(s, "先把约束修成可微，再用二阶方法一步解到底，最后只把该解的变量放进方程。",
             7.14, 4.96, 5.30, .80, 15, WHITE, True)
    add_text(s, "结构性加速来自只更新违反约束的对象与自由度，以及精确 Schur 消元；"
                "二阶求解器是其中的求解层，其默认状态仍在 Smoke 验证中。",
             7.14, 5.82, 5.30, .44, 11, RGBColor(199, 219, 238))

    # ------------------------------------------------------------ 11 SceneProof
    s = base_slide(prs, "SceneProof：proposal 不等于 commit",
                   "每个候选都是一笔事务；只有局部通过且全局非劣才提交", page())
    steps = [("候选", "LM / rule", BLUE), ("重建", "true mesh", PURPLE), ("局部门", "contact / COM", ORANGE),
             ("全局门", "family non-regression", TEAL), ("裁决", "commit / restore", GREEN)]
    xs = [.55, 3.03, 5.51, 7.99, 10.47]
    for i, (title, detail, color) in enumerate(steps):
        add_box(s, xs[i], 1.52, 2.12, 1.18, WHITE, color, width=1.5)
        add_text(s, title, xs[i] + .18, 1.75, 1.76, .30, 16, color, True, PP_ALIGN.CENTER)
        add_text(s, detail, xs[i] + .15, 2.18, 1.82, .25, 10.4, MID, False, PP_ALIGN.CENTER)
        if i < 4:
            add_arrow(s, xs[i] + 2.15, 2.12, xs[i + 1] - .05, 2.12, MID, 1.8)
    add_box(s, .68, 3.20, 5.80, 2.54, WHITE, GRID)
    add_text(s, "必须出现的证据", .98, 3.51, 2.6, .34, 18, NAVY, True)
    add_bullets(s, ["真实网格不增加碰撞 / 穿透",
                    "declared support contact + COM / boundary",
                    "plane / semantic / visibility 不明显退化"],
                .98, 4.05, 4.92, 1.30, 13.2)
    add_box(s, 6.78, 3.20, 5.88, 2.54, PALE_TEAL, TEAL)
    add_text(s, "提交规则", 7.10, 3.51, 2.0, .34, 18, TEAL_DARK, True)
    add_text(s, "Accept(\u0394) \u21d4 LocalGates(\u0394) \u2227 \u0394Family \u2265 \u2212\u03b5",
             7.11, 4.02, 5.45, .62, 17, NAVY, True, PP_ALIGN.CENTER, font=MATH)
    add_text(s, "失败 \u2192 restore incumbent\n证据不足 \u2192 unresolved（不伪装成功）",
             7.12, 4.74, 5.12, .72, 14, DARK, True, PP_ALIGN.CENTER)

    # ------------------------------------------ 12 physical macro 的定义
    s = base_slide(prs, "physical macro 是怎么定义的，为什么它可信",
                   "它的价值不在分数高，而在每个测不了的量都被标成测不了，而不是记零分", page())
    add_box(s, .62, 1.36, 6.02, 2.02, WHITE, TEAL, width=1.6)
    add_text(s, "打分核：误差达到容差即归零", .92, 1.60, 4.4, .30, 16, TEAL_DARK, True)
    add_text(s, "score = max( 0, 1 \u2212 error / tolerance )", 1.00, 1.96, 5.26, .52, 16,
             NAVY, True, PP_ALIGN.CENTER, font=MATH)
    add_text(s, "支撑项由三个等权加数合成：接触间隙（容差 5 cm）、外扩误差（5 cm）、"
                "落脚面重叠比（目标 0.9）。碰撞项取每个物体最坏的那一对重叠体积比（容差 5%）。",
             .96, 2.46, 5.40, .84, 12.2, DARK)
    add_box(s, 6.82, 1.36, 5.88, 2.02, PALE_BLUE, BLUE, width=1.6)
    add_text(s, "两级聚合", 7.12, 1.60, 3.0, .30, 16, BLUE, True)
    add_text(s, "macro = mean(四族分数)\ncritical = min(四族分数)", 7.14, 1.92, 5.30, .58, 15,
             NAVY, True, PP_ALIGN.CENTER, font=MATH)
    add_text(s, "先在族内按物体平均，再按场景等权平均。macro 用平均是为了反映整体可用性；"
                "同时报 critical 取最小值，防止一个族的崩塌被其他三族的高分掩盖。",
             7.14, 2.54, 5.30, .80, 12.2, DARK)
    add_text(s, "三条让它不自欺的设计", .92, 3.62, 5.0, .32, 18, NAVY, True)
    honest = [("不可测量就弃权，不给零分。",
               "房间的空间范围在任何产物里都没有被表示，所以物体是否在房间内这个问题不是答得好坏，"
               "而是答不了。曾试图用墙面反推地板参考面，实测推出的面落在地板原点下方 2.87\u20133.15 m"
               "（预期 +0.02 m），假设被证伪后代码删除，而不是调参保留。"),
              ("碰撞分数被声明为下界。",
               "碰撞用包围盒实心柱计算，它包含真实网格，所以报出的重叠对是真实穿透的完备超集。"
               "分数因此是更精细几何下的下界，这句话写进产物字段，而不是口头说明。"),
              ("弃权会改变估计目标，因此不可跨运行比较。",
               "被弃权的项恰恰是最可能低分的那些，剩下的平均回答的已是另一个问题。"
               "产物里因此带有 estimand_changed_by_abstention 标记。")]
    y = 4.06
    for i, (head, body) in enumerate(honest):
        add_box(s, .62, y, 12.08, .68, PALE_TEAL if i == 0 else WHITE, TEAL, width=1.2)
        add_text(s, head, .88, y + .07, 3.60, .26, 12.6, TEAL_DARK, True)
        add_text(s, body, 4.56, y + .05, 7.94, .60, 10.6, DARK)
        y += .78

    # --------------------------------------------------------------- 13 质量
    s = base_slide(prs, "质量结果：V3 找回更多，V5 让场景更物理",
                   "Paper30 \u00b7 Primary objects \u00b7 visible mask \u2265 8,000 px \u00b7 GT 只用于评测", page())
    data = [["版本", "Primary recovery", "Primary parent", "Physical macro", "它解决什么"],
            ["Imaginarium V1", "89.49%", "89.32%", "52.98%", "原始基线"],
            ["Lumenarium V3", "91.40%", "87.80%", "52.14%", "支撑感知恢复"],
            ["V4 DeepSearch", "88.22%", "80.14%", "54.58%", "检索提速，pose 下降"],
            ["Lumenarium V5-fast", "88.22%", "80.14%", "62.10%", "速度 + 物理主版本"]]
    add_table(s, data, .62, 1.43, 12.08, 2.66, [2.4, 1.8, 1.7, 1.65, 2.4], 12.4, {2: PALE_ORANGE, 4: PALE_TEAL})
    add_metric(s, .68, 4.48, 3.70, 1.55, "91.40%", "V3 recovery", "结构化支撑帮助找回 Primary 对象", ORANGE)
    add_metric(s, 4.78, 4.48, 3.70, 1.55, "+7.52 pp", "V5 physical macro", "54.58% V4 \u2192 62.10% V5", GREEN)
    add_box(s, 8.88, 4.48, 3.78, 1.55, PALE_BLUE, BLUE)
    add_text(s, "为什么这 7.52 pp 能归因给 S4", 9.17, 4.72, 3.4, .30, 14, BLUE, True)
    add_text(s, "V4 与 V5 的 recovery / parent 完全相同（88.22% / 80.14%），上游工作点被冻结，"
                "因此差异只能来自 S4 的优化与证明层。",
             9.17, 5.10, 3.30, .84, 11.4, DARK)

    # --------------------------------------------------------------- 14 速度
    s = base_slide(prs, "速度结果：全链路从约 23.8 分钟降到 13.83 分钟",
                   "V1 为历史恢复估算；V4 为 S0\u2013S3 实测 + legacy S4；V5 为 Paper30 冷启动实测", page())
    data = [["版本", "端到端 / scene", "S4 / scene", "相对 V1", "主要变化"],
            ["Imaginarium V1", "\u224823.8 min", "677.770 s", "1.00\u00d7", "legacy retrieval + SA-5000"],
            ["V4 DeepSearch", "\u224821.9 min", "677.770 s", "\u22481.09\u00d7", "检索提速"],
            ["Lumenarium V5-fast", "13.83 min", "192.930 s", "\u22481.72\u00d7", "SceneLM + SceneProof"]]
    add_table(s, data, .62, 1.43, 12.08, 2.38, [2.30, 1.72, 1.55, 1.30, 3.03], 12.4, {3: PALE_TEAL})
    add_metric(s, .70, 4.24, 3.72, 1.62, "3.513\u00d7", "S4 speedup", "677.770 s \u2192 192.930 s", TEAL)
    add_metric(s, 4.80, 4.24, 3.72, 1.62, "13.83 min", "V5 S0\u2013S4", "829.879 s / scene measured", BLUE)
    add_box(s, 8.90, 4.24, 3.72, 1.62, PALE_ORANGE, ORANGE)
    add_text(s, "下一个瓶颈很明确", 9.19, 4.46, 2.8, .28, 15, ORANGE, True)
    add_text(s, "S1 = 443.0 s，占冷启动 S0\u2013S3 的 69.6%，而 Gemini 当前有效并发只有 1。",
             9.19, 4.80, 3.24, .56, 11.2, DARK)
    add_text(s, "开到 8 路并发并合并请求，单场景约省 200 s，端到端落到约 10 分钟。"
                "这是容量规划估算，不是实测值。",
             9.19, 5.40, 3.24, .70, 11.2, ORANGE, True)
    add_box(s, .62, 6.06, 12.08, .92, NAVY, NAVY)
    add_text(s, "长期成本视角", .92, 6.20, 2.2, .26, 13, TEAL, True)
    add_text(s, "一个约 30 个资产的场景，现在要 10\u201315 分钟。但这个数字里有两种性质完全不同的成本："
                "S1 的 443 s 是被限流迫出的串行等待，不是算力不够，它随并发额度与算力线性缓解；"
                "而 S4 的 3.513\u00d7 是算法层的收益，已经落袋，不依赖未来的硬件。",
             3.10, 6.18, 9.32, .52, 11.6, WHITE, False)
    add_text(s, "所以随着算力与并发上来，单场景成本会继续往下走，而已经拿到的算法收益不会退回去。",
             3.10, 6.66, 9.32, .24, 11.6, TEAL, True)

    # ------------------------------------------------------------ 18/19 视觉
    s = base_slide(prs, "同一客厅：从" + q("重建出来") + "到" + q("整体可用"),
                   "Input / Imaginarium V1 / support-aware V3 / final Lumenarium V5", page())
    add_four_up(s, living, ["Input", "V1 \u00b7 Imaginarium", "V3 \u00b7 Support-aware", "V5 \u00b7 Lumenarium"])
    add_box(s, .70, 5.82, 11.96, .58, PALE_TEAL, TEAL)
    add_text(s, "V5 更完整地恢复大件与软装，并把系统目标从单体 pose 转向场景级可用性。",
             .96, 5.98, 11.42, .28, 14.2, TEAL_DARK, True, PP_ALIGN.CENTER)

    s = base_slide(prs, "同一办公室：更快检索，同时保持关键资产覆盖",
                   "视觉结果不是 GT；它展示的是产品级完整性、关系和冲突 trade-off", page())
    add_four_up(s, office, ["Input", "V1 \u00b7 Imaginarium", "V3 \u00b7 Support-aware", "V5 \u00b7 Lumenarium"])
    add_box(s, .70, 5.82, 11.96, .58, PALE_ORANGE, ORANGE)
    add_text(s, "DeepSearch 提升速度和资产覆盖，但单体 rotation / translation 仍需后续校准。",
             .96, 5.98, 11.42, .28, 14.2, DARK, True, PP_ALIGN.CENTER)

    # --------------------------------------------------------------- 17 服务
    s = base_slide(prs, "不止是论文代码：已经变成双 A10 端到端服务",
                   "新图完整运行 S0\u2013S4；相同输入跨 Fast / Medium / Best 复用冻结缓存", page())
    add_box(s, .62, 1.43, 3.14, 4.82, WHITE, GRID)
    add_text(s, "技术美术入口", .95, 1.76, 2.3, .34, 19, NAVY, True)
    add_bullets(s, ["任意尺寸 PNG / JPEG", "S0\u2013S4 实时进度", "可编辑 placement + render", "证书 + unresolved + ZIP"],
                .96, 2.38, 2.35, 2.40, 14)
    add_box(s, 4.05, 1.43, 4.18, 4.82, PALE_BLUE, BLUE)
    add_text(s, "可靠调度", 4.38, 1.76, 2.0, .34, 19, BLUE, True)
    add_bullets(s, [("双 A10：", "原子 claim，避免重复领取"), ("恢复：", "回收死亡 claim；单场失败不中断"),
                    ("重试：", "HTTP / S4 缺失有限重试"), ("缓存：", "复用 S0\u2013S3 / Fix61")],
                4.39, 2.35, 3.30, 2.78, 12.2, bullet_color=BLUE)
    add_box(s, 8.53, 1.43, 4.14, 4.82, PALE_TEAL, TEAL)
    add_text(s, "三种模式", 8.86, 1.76, 2.0, .34, 19, TEAL_DARK, True)
    add_bullets(s, [("Fast：", "冻结 Fix61，论文定量"), ("Medium：", "visual-safe 展示清理"),
                    ("Best：", "全对象真实支撑审计与事务化 first-contact 掉落")],
                8.87, 2.35, 3.25, 2.48, 12.2, bullet_color=TEAL)
    add_box(s, 8.86, 5.05, 3.25, .56, NAVY, NAVY)
    add_text(s, "embedding.lightart.qq.com", 9.00, 5.21, 2.97, .24, 11.5, WHITE, True, PP_ALIGN.CENTER)

    # ------------------------------------------------------------- 18 自评
    s = base_slide(prs, "我给这两个月打 8.5 分", "扣掉的 1.5 分，是我知道该做、但还没做完的事", page())
    add_box(s, .62, 1.38, 6.02, 4.96, PALE_TEAL, TEAL, width=1.7)
    add_text(s, "8.5 分从哪来", .92, 1.68, 3.0, .34, 20, TEAL_DARK, True)
    add_bullets(s, [("闭环了，不是停在实验：", "算法、评测、服务同时落地，技术美术能直接用。"),
                    ("因果说得清：", "冻结上游后对比，+7.52 pp 能指到具体模块，"
                                    "而不是笼统地说整体变好了。"),
                    ("诚实是设计出来的：", "pose 下降写在报告里；证据不足标 unresolved；"
                                          "测不了的族弃权而不是记零分。"),
                    ("可复现：", "Paper30 冻结产物、provenance、失败即回滚。")],
                .93, 2.28, 5.42, 3.90, 13, bullet_color=TEAL)
    add_box(s, 6.82, 1.38, 5.88, 4.96, PALE_ORANGE, ORANGE, width=1.7)
    add_text(s, "扣掉的 1.5 分", 7.12, 1.68, 3.0, .34, 20, ORANGE, True)
    add_bullets(s, [("pose 工作点没修回来：", "rotation AUC 从 V1 的 48.13% 到现在 31.38%。"
                                            "我定位到它发生在 DeepSearch 之后、SceneProof 之前，"
                                            "但没能在实习内修复。"),
                    ("二阶求解器还没转正：", "关系图、matrix-free LM、Schur、证书都已实现，"
                                            "默认仍是 OFF，等 Smoke 验证。架构完成不等于上线。"),
                    ("最大的时间瓶颈没动：", "S1 占冷启动 69.6%，我一直知道，但优先做了 S4。"),
                    ("评测工具做得比修复多：", "有一段时间连续多轮只增审计不改 pipeline，"
                                              "边际收益衰减时我应该更早收手。")],
                7.13, 2.28, 5.30, 3.90, 12.6, bullet_color=ORANGE)

    # -------------------------------------------------- 19 瓶颈（两类分开）
    s = base_slide(prs, "剩下的问题：哪些是行业硬骨头，哪些是我下周就能动的",
                   "把暂时无解和还没做分开，才知道力气该往哪使", page())
    add_box(s, .62, 1.38, 6.02, 4.96, PALE_RED, RED, width=1.6)
    add_text(s, "公认硬性问题（短期无解）", .92, 1.68, 4.4, .34, 19, RED, True)
    add_bullets(s, [("VLM 质量：", "场景图与语义判断的错误会向下级联；"
                                  "换更强模型只抬高上限，不改变机制。"),
                    ("资产库有限：", "2,043 个资产，检索必然是近似匹配；"
                                    "库里没有的东西，再好的算法也变不出来。"),
                    ("单图尺度歧义：", "从一张图恢复真实尺寸本身欠定。最刺眼的一类错误是凭空长出"
                                      "1.9 m 书架：观测框只剩几厘米碎片时缩放退回 [1,1,1]，"
                                      "而这个看似中性的默认值，实际是在断言此物恰好就是"
                                      "资产授权的那么大。")],
                .93, 2.28, 5.42, 3.90, 12.4, bullet_color=RED)
    add_box(s, 6.82, 1.38, 5.88, 4.96, PALE_TEAL, TEAL, width=1.6)
    add_text(s, "眼前就能解决的", 7.12, 1.68, 4.0, .34, 19, TEAL_DARK, True)
    add_bullets(s, [("S1 并发：", "Gemini 从有效并发 1 提到 8 并合并请求，"
                                 "单场景约省 200 s，端到端 13.83 \u2192 约 10 分钟。"),
                    ("二阶求解器转正：", "Smoke1 / Smoke5 通过后成为主路径。"),
                    ("资产尺寸先验：", "资产库授权尺寸与网格的一致率已实测 91.4%，"
                                      "可以用它给缩放兜底，替掉 [1,1,1] 这个危险默认值。"),
                    ("pose 校准：", "seed-locked S2 ablation，定位 DeepSearch 的姿态代价。")],
                7.13, 2.28, 5.30, 3.90, 12.4, bullet_color=TEAL)

    # ------------------------------------------------------- 20 拓展与思考
    s = base_slide(prs, "如果继续做下去，我认为真正的方向是什么",
                   "这两个月最有迁移价值的，不是某个模块，而是一种做法", page())
    add_box(s, .62, 1.38, 12.08, 1.50, NAVY, NAVY)
    add_text(s, "任何生成式系统都有同一个问题：模型的输出是提议，不是结论。",
             1.00, 1.62, 11.3, .40, 21, WHITE, True, PP_ALIGN.CENTER)
    add_text(s, "我们在三维场景里给这件事装了一层证书。同样的结构可以装到别的地方去。",
             1.00, 2.14, 11.3, .34, 15, TEAL, True, PP_ALIGN.CENTER)
    cards = [("近期", "把证书变成可微的",
              "现在的证书是事后裁决：先改、再验、不行就回滚。如果把接触与稳定性写成可微项，"
              "就能在优化中直接朝可证明的方向走，而不是靠事后筛。", BLUE),
             ("中期", "闭环到 DCC",
              "闭环的前提是先把场景完整交出去。现在导出的 .glb 位姿与资产引用都是对的，"
              "但多部件资产只装配了第一个部件，所以它适合看与验证，还不能当成品交付。"
              "把这一步做实之后，技术美术在 Blender / UE 里的手工修正就是最高质量的监督信号，"
              "回写成约束，系统会越用越准。", PURPLE),
             ("远期", "同一套方法迁移到 agent",
              "LLM agent 调用工具时面对的是同一个问题：proposal 不等于 commit，"
              "需要局部门加全局非劣，以及证据不足时敢说 unresolved。", TEAL)]
    for i, (tag, head, body, color) in enumerate(cards):
        x = .62 + i * 4.09
        add_box(s, x, 3.12, 3.86, 3.22, WHITE, color, width=1.6)
        add_text(s, tag, x + .28, 3.36, 1.2, .28, 13, color, True)
        add_text(s, head, x + .28, 3.72, 3.30, .64, 17, NAVY, True)
        add_text(s, body, x + .28, 4.52, 3.30, 1.62, 12.4, DARK)

    # --------------------------------------------------------------- 21 影响
    s = base_slide(prs, "这两个月留下的，不只是一个更高的分数", "算法、证据、工具链和产品入口同时闭环", page())
    cards = [("研究", "两项方法创新", "support-aware reconstruction\nSceneLM + SceneProof", BLUE),
             ("工程", "可复现评测", "Paper30 / 8000px+\nprovenance / rollback", ORANGE),
             ("产品", "双 A10 服务", "Web / API / cache\nFast / Medium / Best", TEAL)]
    for i, (tag, head, body, color) in enumerate(cards):
        x = .70 + i * 4.18
        add_box(s, x, 1.55, 3.78, 3.50, WHITE, color, width=1.7)
        add_text(s, tag, x + .30, 1.88, 1.0, .28, 13, color, True)
        add_text(s, head, x + .30, 2.41, 3.02, .38, 21, NAVY, True)
        add_text(s, body, x + .30, 3.17, 3.02, 1.05, 15, DARK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_box(s, 1.10, 5.48, 11.14, .72, NAVY, NAVY)
    add_text(s, "把单图 3D 场景生成，从研究原型推进成更快、可证明、"
                "能被技术美术直接使用的系统。",
             1.38, 5.70, 10.58, .31, 17, WHITE, True, PP_ALIGN.CENTER)

    # ------------------------------------------------------------- 22 致谢
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.13))
    set_fill(band, TEAL); band.line.fill.background()
    bottom = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.37), prs.slide_width, Inches(.13))
    set_fill(bottom, TEAL); bottom.line.fill.background()
    add_text(s, "谢    谢", .90, .70, 6.4, 1.34, 72, WHITE, True)
    add_text(s, "T H A N K   Y O U", .96, 2.14, 6.0, .40, 15, TEAL, True)
    add_box(s, .90, 2.82, 6.30, 2.66, NAVY2, TEAL, width=1.6)
    add_text(s, "两个月前，我只会说这个场景看起来不太对。", 1.20, 3.10, 5.74, .42, 17, WHITE, True)
    add_text(s, "现在我能说清它错在哪一项、差多少、为什么，以及凭什么敢改。",
             1.20, 3.62, 5.74, .78, 17, TEAL, True)
    add_text(s, "这个转变不是我一个人完成的。", 1.20, 4.54, 5.74, .34, 15, WHITE, False)
    add_text(s, "谢谢每一次把我的差不多可以了打回来的 review。",
             1.20, 4.96, 5.74, .38, 15, RGBColor(202, 220, 239), True)
    add_box(s, 7.56, 2.82, 4.88, 2.66, NAVY2, TEAL, width=1.6)
    add_text(s, "特别感谢", 7.86, 3.10, 3.0, .34, 17, TEAL, True)
    add_text(s, "Calvin Gu", 7.86, 3.56, 4.3, .50, 26, WHITE, True)
    add_text(s, "谢谢你在我想走捷径的时候，一直问我怎么证明。这四个字改变了我做研究的方式。",
             7.86, 4.20, 4.32, .96, 13.5, RGBColor(214, 229, 245))
    add_text(s, "以及 Lumenarium 与光子团队的每一位同事：谢谢你们把一个实习生的问题，"
                "当成真问题来回答。",
             7.86, 5.24, 4.32, .90, 13, RGBColor(202, 220, 239))
    add_runs(s, [("完整性  ", TEAL, True), ("\u2192 速度  ", WHITE, False), ("\u2192 低穿模  ", WHITE, False),
                 ("\u2192 关系可信  ", WHITE, False), ("\u2192 可交付", WHITE, False)],
             .92, 5.72, 11.50, .44, 19, MSO_ANCHOR.MIDDLE, PP_ALIGN.CENTER)
    add_text(s, "Questions & Discussion", .92, 6.40, 6.0, .40, 22, WHITE, True)
    link_text(add_text(s, "Hansen Zhu \u00b7 Mentor: Calvin Gu \u00b7 Demo: " + DEMO_VIDEO_URL,
                       .94, 6.92, 11.44, .28, 12, RGBColor(202, 220, 239)))

    reorder(prs)
    attach_notes(prs)  # 必须在 reorder 之后：备注按放映顺序对齐讲稿
    prs.save(OUT)
    return OUT


# 讲稿是单一真源：同一份内容既注入 PPT 的备注栏（放映视图里看得见），
# 也导出成 markdown 逐页稿。每页三样东西——秒数、可以直接念出来的话、
# 以及只给自己看的动作提示。幻灯片上不再写动作提示：投到屏幕上的每一句，
# 都应该是听众需要的信息。
SEGMENTS = (
    ("开场与问题", 1, 5),
    ("系统全貌与主线", 6, 7),
    ("录屏与发起演示", 8, 9),
    ("支撑线与方法", 10, 14),
    ("指标与结果", 15, 17),
    ("产出与现场验收", 18, 20),
    ("自评与展望", 21, 25),
)

SPEECH = (
    (10, "封面",
     ("各位好。这两个月我把一个跑得通的研究原型，"
      "推成了更快、能自己证明结果、可以直接交付的系统。",),
     ("一句话讲完就翻页，不要在封面停留。",)),
    (21, "技术美术的一天",
     ("先讲人。技术美术要的不是一张看着像的图，是一个导进引擎就能开工的场景。",
      "可看起来对的场景常常不能用：杯子浮空两厘米，动画一动就露；"
      "窗帘没检索出来，整个镜头要重做。",
      "所以这两个月我在回答：怎么让系统只提交它能证明的修改。"),
     ("三个例子只挑两个说：杯子和窗帘。",
      "落到右边那句大字上再翻页，这是全场的主线句。")),
    (12, "交付优先级",
     ("先定义什么叫好：资产完整、速度、不穿模、关系成立。取舍也写明了"
      "\u2014\u2014单体朝向可以稍作牺牲，少量悬空宁可报 unresolved。",),
     ("念四条排序，右边取舍栏一句话带过，不展开。",)),
    (36, "继承的 Imaginarium",
     ("说明一下我接手时的样子：先粗排、再用 VLM 比对纹理和尺寸精排，最后把三维框对回二维观测。",
      "开销集中在两处。一个是限流之下 API 调用全在串行等待；更关键的是支撑关系靠文字判断"
      "\u2014\u2014VLM 会把柜顶的箱子写成放在地上，物理一跑就穿透地面。",
      "所以最后一步要用五千步退火去兜：前面每步不确定，后面只能用次数换收敛概率。"
      "这就是我从上游入手的原因。"),
     ("收在「前面不确定、后面就得兜」这句因果上，这是整段的目的。",
      "底部取舍表第三行是刻意不做 Flux 微调：承认 Top-1 从 48.6% 到 68.7% 的收益，"
      "但那是对某一批资产过拟合，而我们要交付能换库的系统。被问到再展开。")),
    (13, "V1 的四类痛点",
     ("这是当时的输出。能生成，但四类问题挡在交付前面：资产缺失、太慢、穿模、关系不成立。",
      "后面所有工作都是冲着这四条去的。"),
     ()),
    (60, "全链路 S0\u2013S4",
     ("这是现在的完整链路，五个阶段。我按每一步吃什么、吐什么过一遍，"
      "后面所有的痛点和改进都挂在这上面。",
      "S0 吃一张概念图，吐单目深度、相机内参和点云\u2014\u2014它给后面所有几何提供尺度。",
      "S1 吃图和深度，吐实例分割、有向包围盒，和一张场景图，也就是谁被谁支撑。",
      "S2 吃物体切图和包围盒，去资产库里找最像的模型，吐 Top-K 候选。",
      "S3 吃候选资产、包围盒和 mask，用几何与视觉双通道打分选定资产，"
      "再把三维框对齐回二维观测，吐每个物体的六自由度位姿与缩放。",
      "S4 吃全部位姿加支撑树，吐无碰撞布局、渲染图，和一份几何证书。",
      "我重做的是后三个：检索、姿态、以及最后这一步的优化与提交。"),
     ("这一页是「带过但要讲清」：五个阶段的接口都说出口，不要只指着方框说五个阶段。",
      "被问到 S1 细节：GroundingDINO 或 SAM3 出 mask，DepthAnything V2 出深度，"
      "mask 内像素按深度反投影成点云再拟合 OBB。")),
    (7, "版本主线",
     ("这是两个月的主线：V1 到 V3、V4、V5，最后落成一个在线服务。",),
     ("手指着五个节点扫一遍就行，不要逐版解释，后面的页会讲。",)),
    (22, "三档实机录屏",
     ("知道链路长什么样之后，看一段它跑完的录屏。同一张输入跑三档模式，剪辑过、也加了速。",
      "最重要一句：三档不是三个模型，而是同一条流水线上三种提交策略。"
      "越往右提交越保守、要求的证据越多；重建结果共用，换档只重跑最后一步。"),
     ("放视频，少说话，让画面自己走。",
      "主动说明是剪辑并加速的\u2014\u2014承认的成本很低，被人看出来的代价很高。",
      "封面那块牌子和底部链接都可点，直接跳 B 站；若已把 mp4 放进仓库根目录，则是原地播放。")),
    (25, "实机演示",
     ("录屏是剪过的，所以现在真跑一次。这是跑在两块 A10 上的在线服务，"
      "给它一张参考图，拿到能直接在 Blender 里打开的场景。",
      "我上传左边这张，选 V5-demo，随机种子固定。今天不用缓存，从头全量跑，"
      "大概十几分钟。它在后台跑，我们接着讲。"),
     ("先把图上传、选好档、点 Generate，让进度条动起来，再开口说话。",
      "这一页只放输入图。输出此刻还不存在，跑完直接切浏览器看，不要往幻灯片里补图。",
      "网络不通就说「刚才那段录屏就是它的完整过程」，往下走，不要在台上 debug。")),
    (56, "V3 的 support-aware 怎么做",
     ("回到刚才那个错误。先分清两个环节：检测是谁干的，关系是谁判的。"
      "检测是视觉的\u2014\u2014GroundingDINO 或 SAM3 出框和 mask，这是模型在看像素；"
      "但「谁在谁上面」这个关系，在继承的系统里是 GPT 用文字回答的\u2014\u2014"
      "它会把柜顶的箱子说成放在地上。",
      "V3 的 support-aware，也就是堆叠感知，就是把这个文字判断换成纯几何的："
      "不看任何模型输出，只看 mask 和 OBB 的几何关系。",
      "做法三步。第一步检测堆叠对：两个物体的水平投影重叠超过一定比例、"
      "竖直间隙又落在容差之内，就认定上下关系。所有阈值都由物体自身尺寸算出来，没有硬编码。",
      "第二步是抬升与冻结：先算出上层该落在哪个高度，退火期间把它从被扰动的集合里排除，"
      "退火结束后再按下层顶面精确落位。",
      "第三步沿支撑链递归：父物体一动，它上面所有后代跟着走。",
      "这条线在 V5 完整保留，今天还在跑；S3 侧一个开关都没关。"),
     ("这一页是重中之重，不要压。三步各一句，讲完停半秒再翻页。",
      "被问「几何是不是 DINO 判的」：不是。DINOv2 只在 S2 检索做特征粗排，GroundingDINO/SAM3 "
      "只在 S1 做检测分割，它们都不判关系。V3 起关系是纯几何算的：S1 侧有 mask 底部接触像素的"
      "重叠检测，S3 侧有 OBB 水平投影重叠加竖直间隙的堆叠检测，两个都是几何算法。",
      "被问到为什么还信 VLM 的高层关系：错基本只错在落地层，高处的父子它判得对，所以链条可用。")),
    (58, "从关系到 SceneLM",
     ("但 support-aware 只回答了「谁在谁上面」。判对之后该怎么摆，V3 还是交给退火去碰运气"
      "\u2014\u2014它没办法说明为什么该动。",
      "所以 V5 换了思路：把这些关系直接编译成一个可以求解的目标函数，这就是 SceneLM。",
      "它把场景要满足的东西显式写成一个式子：不能互相穿透、该接触的要接触、靠墙的要贴墙，"
      "再加语义偏好，比如椅子该朝着桌子。",
      "重点是这张权重表怎么读。几何硬约束的权重是语义偏好的四倍\u2014\u2014好看要让位于成立。"
      "warm-start 只给 0.01，意思是上游初值只是参考，不是要保护的东西。",
      "权衡不再隐含在谁的调参手感里，而是显式写下来、可以被质疑。"),
     ("开头那句转折是全场的枢纽：从「判对关系」到「按关系求解」。",
      "方法段从这里开始不可压。先念总式，再讲权重表。")),
    (52, "为什么敢用二阶",
     ("这一页要回答两个问题：**要优化的到底是什么、为什么我们能用二阶法**。",
      "**要优化的 x**：上一页 SceneLM 给了一个 E(x)，这里的 x 就是每个物体的位姿\u2014\u2014"
      "三维位置加朝向。我们想让它越小越好。",
      "**怎么把约束变成数字**：每条约束都化成" + q("差几厘米 / 几度") + "，"
      "比如穿透了 3 厘米就记 r = 0.03。所有约束堆成一个向量 r(x)，"
      "那 E(x) 就是 r 的模方。我们要找的就是让 r 的模方最小的 x。",
      "**为什么不能直接上 Gauss\u2013Newton**：碰撞约束是「或」的关系\u2014\u2014两物体穿透了，"
      "推开的最省力方向是哪一条？这条最省力轴一旦切换，方向就跳。"
      "二阶法在切点处做线性化，方向一跳也跟着震荡。所以要先把方向修光滑。",
      "**怎么修**：判定某对已穿透，就冻住它最省力的那个分离方向\u2014\u2014"
      "「或」变成一个稳定的标量残差。这时曲面在切点附近是光滑的，"
      "二阶法才稳。**顺序不能颠倒**：先冻方向，再上二阶。",
      "**LM 求解器**：形式是 (J\u1d40J + \u03bb\u00b7diag(J\u1d40J))\u03b4 = \u2212J\u1d40r，"
      "\u03b4 就是这一步要走的位姿增量。**我们从不构造 J 与 Hessian**："
      "J\u1d40J\u00b7v 由两次自动微分得出，外层 PCG 迭代解，"
      "内存与一阶同阶，却拿到二阶的收敛性。**\u03bb 自调**：实际下降 / 预测 \u2265 0.75 就减半、"
      "更像纯二阶；\u2264 0.25 就加倍、退回梯度下降；整步被拒就乘四。",
      "一句话：把约束堆成残差向量，先冻分离方向让曲面光滑，然后用 LM 一步到二次模型的谷底。"),
     ("务必按「要优化 x → 残差向量 → 不光滑所以冻方向 → LM 求解」顺序讲，顺序本身就是论证。",
      "被追问 \u03b4 是什么：这一步要加到 x 上的位姿增量，LM 步通常不缩放（不像梯度下降还要乘学习率）。",
      "被追问为什么不在中间构造 Hessian：场景大了存不下，自动微分只用两次前向+反向，内存不爆。",
      "超时就先砍 \u03bb 那段规则，直接说「λ 会自适应」。")),
    (54, "为什么快 \u00b7 诚实边界",
     ("承上一页\u2014\u2014上一节的 LM 求解式 (J\u1d40J + \u03bb\u00b7diag(J\u1d40J)) \u03b4 = \u2212J\u1d40r "
      "已经能解了，这一页要回答\u201c它为什么快\u201d，以及\u201c什么时候失效\u201d。",
      "\u201c为什么会快，原因不在用了二阶\u201d，而在两件更朴素的事。",
      "\u201c第一\u201d：关系本身决定了要动几个自由度。被支撑的两个，贴墙的两个，法向锚定一个，"
      "完全自由的才三个。LM 只更新真正违反约束的那些物体\u2014\u2014其他物体的位姿根本不动。",
      "\u201c第二\u201d：Schur 消元只消掉有父节点、没有子节点的叶子平移。"
      "对当前这次线性化来说，这是精确消元，不是近似。",
      "\u201c诚实边界\u201d：LM 在接触状态切换的瞬间不可靠\u2014\u2014两物体从分离跳到接触的那一帧，"
      "Gauss\u2013Newton 的线性化不再贴切。这时我们做一件事：把 LM 的当前步投影到"
      "\u201c已审计的责任子空间\u201d（即上一轮 SceneProof 标记为违反约束的那些方向），"
      "在这个低维子空间里做无导数正交搜索\u2014\u2014试几个方向、留下最好的那一步。"
      "它比 LM 慢，但\u201c只在例外情形触发\u201d，不影响日常路径的速度。",
      "\u201c口径纪律\u201d：3.513 倍的加速来自刚说的两点结构性优化，"
      "\u201c不要归因给二阶求解器\u2014\u2014它默认还是关着的，等 Smoke 验证过再说。\u201d"),
     ("最后那句口径必须说出口，这是防止别人替你夸大的唯一办法。",
      "被追问\u201c退化的步会不会很慢\u201d：不会，它只在接触切换那一帧触发几次，"
      "日常的滑动与堆叠仍是 LM 主路径。")),
    (34, "SceneProof",
     ("关系写成约束、也解出来了，还剩最后一个问题：优化器给出的新位置，凭什么就能算数？",
      "这就是 SceneProof。每次改动都当提议，不当结论：先在局部量一遍"
      "\u2014\u2014接触是不是真的量到了、有没有新的穿透；再看全局指标有没有变差。"
      "两个都过才提交，一个不过就整步回滚。",
      "既证明不了、又不能安全回滚的，就明确标成 unresolved，而不是记成成功。"),
     ("五个方框不要逐个念，指着流程把三种结局说清：提交、回滚、unresolved。",)),
    (32, "physical macro 怎么定义",
     ("指标是这么定义的：误差越小分越高，超过容差归零；支撑这项由三部分相加；"
      "总分取平均，同时另报一个取最小值的版本，防止个别严重错误被平均掉。",
      "这页我最想说中间这条。我原本想用墙面反推地板高度，实测发现推出来的面落在地板下方"
      "两米八到三米一，而预期是正两厘米。所以我把这段代码删了，而不是调参数把它救回来。"),
     ("三条设计只详说墙面反推那条，另两条各一句。这页最能体现研究态度。",)),
    (20, "质量结果",
     ("结果分两块看：恢复率上 V3 最好，物理可实现性上 V5 最好\u2014\u2014它们解决的不是同一个问题。",
      "归因是干净的：V4 和 V5 的恢复率与父子关系完全相同，"
      "所以这 7.52 个百分点只可能来自最后一步。"),
     ()),
    (24, "速度结果",
     ("时间从 23.8 分钟降到 13.83 分钟，最后一步单独快了 3.513 倍。",
      "接上前面的伏笔：当初 API 调用贵，是限流导致串行等待。现在第一阶段占冷启动近七成，"
      "并发开到八路估算能省两百秒\u2014\u2014这是容量规划估算，不是实测。"),
     ("「估算不是实测」这句一定要说，否则会被追问数据来源。",
      "被追问「十几分钟算快吗」：里面有两种成本，一种随算力线性下降，"
      "另一种是算法收益、已经落袋。")),
    (6, "客厅对比",
     ("同一张输入，从左到右。最右边是 V5：更完整，而且整体可用。",),
     ()),
    (0, "办公室对比",
     ("这是办公室场景的同一组对比，结论一样。",),
     ("默认跳过\u2014\u2014第 8 页录屏已经覆盖同样结论；超时就先砍这一页。",
      "被追问单体 pose 时才翻回来，答：检索更快、覆盖更全，但单体朝向与位置仍需校准。")),
    (13, "双 A10 服务 + 现场验收",
     ("这套东西现在不只是论文代码：两块 A10、三种模式、缓存和断点恢复都在线上。",
      "刚才那次运行走的就是这条链路，我们切过去看一眼。"),
     ("全量跑通常还没结束，那就展示进度条与阶段灯\u2014\u2014它本身就是「真在跑」的证据。",
      "跑完了就展示场景和状态行里的随机种子。",
      "若结果报 unresolved，主动说：它给出了结果，同时点名了自己无法证明的关系。"
      "宁可报 unresolved，也不把无法验证的东西记成成功。")),
    (28, "自评 8.5 分",
     ("给自己打 8.5 分。左边是做到的：闭环、因果说得清、诚实是设计出来的、可复现。",
      "扣掉的一分半说具体些：姿态指标没修回来，从 48.13% 掉到 31.38%，"
      "我定位到它发生在哪两步之间，但没能在实习内修好。",
      "二阶求解器写完了却默认还关着\u2014\u2014架构完成不等于上线。"
      "还有一段时间我只加审计不改代码，边际收益衰减时应该更早收手。"),
     ("这一页是主动的自我批评，比任何成绩都更能建立信任。语速放慢。",)),
    (30, "两类瓶颈",
     ("剩下的问题分两类。左边短期无解：VLM 的质量、资产库只有两千零四十三个、单图尺度本身欠定。",
      "尺度这条我讲个具体机制：观测框只剩几厘米碎片时，缩放会退回 1。它看着中性，"
      "其实是在断言「这东西恰好就是资产授权的那么大」，于是库把它翻译成一个 1.9 米高的书架。",
      "右边四条是下周就能动的。"),
     ("把「暂时无解」和「还没做」分开，是这一页存在的全部意义。",)),
    (21, "拓展与思考",
     ("如果继续做，我认为最有迁移价值的不是某个模块，而是一种做法：模型的输出是提议，不是结论。",
      "近期把证书变成可微的；中期闭环到 DCC\u2014\u2014前提是先把场景完整交出去；"
      "远期同一套结构可以搬到 agent 上。"),
     ("中期那张卡写了已知缺口：导出的 .glb 位姿与资产引用都对，"
      "但多部件资产目前只装配第一个部件，所以适合看与验证、还不适合当成品。",
      "被问到交付格式时直接用这句回答，不要含糊地说「可以」。")),
    (7, "影响",
     ("研究、工程、产品这三条线是同时闭环的，这也是我自己最看重的部分。",),
     ("与自评页有重叠，不要重复讲成绩。",)),
    (18, "致谢",
     ("最后。两个月前我只会说「这个场景看起来不太对」。"
      "现在我能说清它错在哪一项、差多少、为什么，以及凭什么敢改。",
      "谢谢 Calvin\u2014\u2014谢谢你在我想走捷径的时候，一直问我怎么证明。"),
     ("慢下来，看着人说。说完停住，进入提问。",)),
)


def notes_text(index: int) -> str:
    seconds, name, say, cues = SPEECH[index - 1]
    text = f"[{index:02d}/{TOTAL}] {seconds}s \u00b7 {name}\n\n" + "\n\n".join(say)
    if cues:
        text += "\n\n\u2014 只给自己看 \u2014\n" + "\n".join("\u00b7 " + cue for cue in cues)
    return text


def attach_notes(prs) -> None:
    """把讲稿写进每页备注栏。

    放映视图里备注是可见的，所以现场提示应该待在这里，而不是印在幻灯片上给听众看。
    """
    if len(SPEECH) != TOTAL:
        raise SystemExit(f"SPEECH 有 {len(SPEECH)} 页，但这份 deck 是 {TOTAL} 页")
    for index, slide in enumerate(prs.slides, 1):
        slide.notes_slide.notes_text_frame.text = notes_text(index)


def pace_report() -> str:
    """按语速区间估算讲稿念完要多久，并报出偏差最大的三页。

    语速是一个假设，不是测量：正式汇报的中文语速通常在每秒 4 到 5 个字之间，
    所以这里给的是区间而不是单一数字。只报偏差、不自动改秒数——
    秒数是这份汇报的时间预算，该被迫改的是内容长度。
    """
    # 标 0 秒的页是刻意跳过的（内容只在被追问时用），把它算进偏差会掩盖真正超长的页。
    rows = sorted(
        (abs(sum(len(part) for part in say) / 4.5 - seconds), index, seconds,
         sum(len(part) for part in say))
        for index, (seconds, _, say, _) in enumerate(SPEECH, 1) if seconds
    )
    worst = ", ".join(f"p{index}={chars}字/{seconds}s" for _, index, seconds, chars in rows[-3:])
    chars = sum(len(part) for _, _, say, _ in SPEECH for part in say)
    labelled = sum(row[0] for row in SPEECH)
    return (f"speech {chars} 字 \u2248 {chars / 5:.0f}\u2013{chars / 4:.0f}s，"
            f"标注 {labelled}s；偏差最大：{worst}")


def segment_table() -> list[str]:
    lines = ["| 段 | 页 | 秒数 |", "| --- | --- | --- |"]
    for name, first, last in SEGMENTS:
        seconds = sum(SPEECH[index - 1][0] for index in range(first, last + 1))
        lines.append(f"| {name} | {first}\u2013{last} | {seconds} |")
    lines.append(f"| **合计** | 1\u2013{TOTAL} | **{sum(row[0] for row in SPEECH)}** |")
    return lines


def speech_markdown() -> list[str]:
    lines: list[str] = []
    for index, (seconds, name, say, cues) in enumerate(SPEECH, 1):
        lines += [f"## {index:02d} \u00b7 {name}（{seconds}s）", ""]
        for paragraph in say:
            lines += [paragraph, ""]
        if cues:
            lines.append("**只给自己看**：")
            lines += [f"- {cue}" for cue in cues]
            lines.append("")
    return lines


def write_notes() -> Path:
    L, R = LQ, RQ
    lines = [
        "# Lumenarium 暑期实习总结（完整版）逐页讲稿",
        "",
        f"每页都给了可以直接念出来的话，逐页秒数加总 **{sum(row[0] for row in SPEECH)} 秒 \u2248 11 分钟**。",
        "**它已经不是 10 分钟了**：把 S0\u2013S4 的接口讲清、把 support-aware 与 SceneLM / SceneProof "
        "的因果讲透，需要这么多时间。如果现场卡在 10 分钟，按这个顺序砍，共约 60 秒："
        "第 19 页办公室对比（已标 0 秒）→ 第 24 页影响 → 第 18 页客厅对比 → 第 7 页版本主线 → "
        "第 3 页优先级只念四个词。**第 6 页与第 10\u201314 页一秒都不要砍**。",
        "每页末尾的「只给自己看」是动作与口径提示：**不在幻灯片上，也不用说出来**。",
        "同一份内容已经写进 PPT 的备注栏，放映视图里能直接看到，不必另开这个文件。",
        "",
        "**顺序是这样安排的**：第 2\u20135 页把人的痛点、什么叫好、以及我们继承到的系统讲完，"
        "第 6 页用一页把 S0\u2013S4 五个阶段的接口交代清楚（**带过但不跳过**：每一步吃什么、吐什么"
        "都要说出口，因为后面所有痛点与改进都挂在这条链路上），第 7 页是版本主线，"
        "第 8 页放录屏，第 9 页才现场跑。听众先知道这条链路长什么样，演示才不显得无来由。",
        "",
        "**第 10\u201314 页是全场重心，一秒都不要压**：第 10 页讲 V3 的 support-aware 到底怎么做"
        "（三步：几何检测堆叠对 \u2192 抬升并在退火期间冻结 \u2192 沿支撑链递归），"
        "第 11 页讲清那个枢纽转折\u2014\u2014判对关系之后怎么摆，"
        "V3 还是交给退火碰运气，所以 V5 把关系编译成可求解的目标函数，这就是 SceneLM；"
        "第 14 页再补最后一环：解出来的位置凭什么算数，这就是 SceneProof。",
        "",
        "**今天的演示是固定随机种子的全量跑**：不复用缓存，整条链路重跑，实测端到端 11\u201318 分钟。"
        "运行在第 9 页末发起，此时已过 182 秒；到第 20 页切回浏览器时它才跑了约 5 分钟，"
        "**大概率还没跑完**。那时就展示进度条与阶段灯\u2014\u2014它本身就是「真在跑」的证据；"
        "完整结果会落在提问环节。",
        "",
        *segment_table(),
        "",
        *speech_markdown(),
        "",
        "",
        "",
        "",
        "",
        "## 预期提问与准备",
        "",
        "- **刚才那次是现场真跑的，还是提前准备好的？** 是真跑的，而且没有用缓存，"
        "整条链路重跑了一遍。",
        "  唯一固定的是随机种子。为什么固定：这条链里有 VLM，同一张图跑两次结果会有差异，"
        "  我们实测过最大 0.38 米。固定种子锁住的是本机的随机源，让演示可复现；",
        "  但它锁不住 VLM 服务端\u2014\u2014我们自己实测过两次全量跑，未证明关系的个数从 9 变成 18。",
        "  所以我只承诺" + L + "同一条链路、同一份配置" + R + "，不承诺逐位一致。",
        "- **为什么有的时候几十秒就出图？** 那是同一张图第二次提交，直接复用了已经缓存的重建结果，"
        "  最后一步的优化、证书和渲染仍然每次都真跑。今天为了完整演示，我们没有用这条路。",
        "- **能不能把跑出来的场景直接给我一个 Blender 文件？** 可以，但要说清它现在是什么：",
        "  下载的结果包里是布局与证书数据（每个物体是资产 ID 加一个 4\u00d74 世界矩阵），几何在资产库里；",
        "  在服务端用 `tools/export_blend.py` 装配后可以导出自包含的 `.glb`，位姿与资产引用都正确。",
        "  **已知缺口**：多部件资产目前只装配第一个部件，所以它适合看与验证，还不适合当成品交付。",
        "  这条写在第 23 页中期那张卡上，**不要含糊地答**" + L + "可以" + R + "。",
        "- **一个场景 10\u201315 分钟，这算快吗？** 先承认它还不够快，再拆成本："
        "  S1 的 443 s 是限流导致的串行等待，属于工程约束，随并发额度与算力线性缓解，",
        "  8 路并发就能省约 200 s；S4 的 3.513x 是算法约束上的突破，已经落袋且不依赖硬件。",
        "  两类成本的下降方式不同，所以长期看单场景成本是持续下行的，",
        "  而且算法侧的收益不会因为换机器而失效。",
        "- **为什么不做 Flux 微调？收益不是很大吗？** 承认收益（Top-1 48.6% \u2192 68.7%），",
        "  但那是把生成图的分布对齐到某一批资产上。一旦换应用场景、换资产库就要重新微调，",
        "  而我们要交付的是能换库的系统。所以选择在检索与优化侧提升鲁棒性，",
        "  而不是在输入侧过拟合。这也解释了我们 pose 指标为何低于论文：是取舍，不是遗漏。",
        "- **为什么 pose 反而变差了？** 定位在 DeepSearch 之后、SceneProof 之前；",
        "  V4 与 V5 工作点相同可证明不是 S4 造成的；已有 seed-locked ablation 方案。",
        "- **physical macro 62% 算高吗？** 它是绝对可实现性的下界，不是相对排名；",
        "  碰撞族被声明为下界，弃权项会改变 estimand，所以只与同口径运行比较。",
        "- **二阶方法为什么还没上线？** 架构与证书已实现，默认 OFF，等 Smoke1 / Smoke5；",
        "  我不把未验证的东西写进主表。",
        "- **这套东西能推广吗？** 见第 23 页第三张卡。",
    ]
    NOTES.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return NOTES


if __name__ == "__main__":
    print(build())
    print(write_notes())
    print(pace_report())
